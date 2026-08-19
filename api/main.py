"""
The web app: an HTML UI plus the JSON API behind it.

Storage is chosen by environment (see `pptx_formatter.config`), so the same
endpoints run against the local filesystem or Supabase without changing.

    uvicorn api.main:app --reload

The UI is at `/`, interactive API docs at `/docs`.

**No authentication.** Every caller has full access to every client's
banked material, so run this on localhost, or put an authenticating proxy
in front of it before it goes anywhere else. When the Supabase backend is
configured it uses the service-role key, which bypasses row-level security
by design - that key stays server-side and is never sent to the browser.
"""
from __future__ import annotations

import json
import shutil
import sys
import tempfile
import uuid
from pathlib import Path

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation                                       # noqa: E402

from pptx_formatter import archetypes                               # noqa: E402
from pptx_formatter.classifier import classify_deck                 # noqa: E402
from pptx_formatter.config import (                                 # noqa: E402
    make_bank, make_job_store, settings_from_env,
)
from pptx_formatter.pipeline import (                               # noqa: E402
    format_with_banked_spec, ingest_master, run_pipeline,
)
from pptx_formatter.style_spec import StyleSpec                     # noqa: E402

app = FastAPI(
    title="PPTX Formatting Tool",
    description=(
        "Submit a designer's master slide, then reformat rough decks to match it."
    ),
    version="1.1.0",
)

STATIC_DIR = Path(__file__).resolve().parent / "static"
WORK_DIR = Path(tempfile.gettempdir()) / "pptx_formatter_uploads"
WORK_DIR.mkdir(exist_ok=True)

PPTX_MEDIA_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


def _settings():
    return settings_from_env()


def _bank():
    try:
        return make_bank(_settings())
    except RuntimeError as exc:
        # Misconfigured Supabase should read as a setup problem, not a crash.
        raise HTTPException(503, str(exc)) from exc


def _jobs():
    try:
        return make_job_store(_settings())
    except RuntimeError as exc:
        raise HTTPException(503, str(exc)) from exc


def _save_upload(upload: UploadFile, suffix: str) -> Path:
    if not (upload.filename or "").lower().endswith(".pptx"):
        raise HTTPException(400, f"Expected a .pptx file, got {upload.filename!r}")
    destination = WORK_DIR / f"{uuid.uuid4().hex}_{suffix}"
    with destination.open("wb") as handle:
        shutil.copyfileobj(upload.file, handle)
    return destination


def _parse_overrides(raw: str | None) -> dict:
    """Accepts `{"4": "quote"}`, keyed by 1-based slide number."""
    if not raw:
        return {}
    try:
        data = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise HTTPException(400, f"overrides must be JSON: {exc}") from exc

    overrides = {}
    for key, value in data.items():
        if value not in archetypes.ALL_ARCHETYPES:
            raise HTTPException(400, f"unknown archetype {value!r}")
        overrides[int(key) - 1] = value
    return overrides


# --- UI -------------------------------------------------------------------

@app.get("/", response_class=HTMLResponse, include_in_schema=False)
def index():
    page = STATIC_DIR / "index.html"
    if not page.exists():                       # pragma: no cover
        raise HTTPException(500, "UI asset is missing")
    return HTMLResponse(page.read_text(encoding="utf-8"))


# --- meta -----------------------------------------------------------------

@app.get("/api/health")
def health():
    """Backend status and the archetype vocabulary the UI renders."""
    return {
        "status": "ok",
        "storage": _settings().describe(),
        "archetypes": [
            {"value": value, "label": archetypes.label_for(value)}
            for value in archetypes.ALL_ARCHETYPES
        ],
    }


# --- stage 1 --------------------------------------------------------------

@app.post("/api/extract")
async def extract(
    master: UploadFile = File(..., description="The designer's master .pptx"),
    client: str | None = Form(None),
    project: str | None = Form(None),
):
    """
    Read a submitted master into a Style Spec and archive it.

    The master may be empty or a fully populated deck; its content slides
    are ignored either way.
    """
    path = _save_upload(master, "master.pptx")
    try:
        spec, entry_id = ingest_master(
            path, bank=_bank(), client=client, project=project
        )
    except HTTPException:
        raise
    except Exception as exc:                                  # noqa: BLE001
        raise HTTPException(422, f"Could not extract a Style Spec: {exc}") from exc

    return JSONResponse({
        "bank_entry": entry_id,
        "summary": _spec_summary(spec),
        "style_spec": spec.to_dict(),
    })


def _spec_summary(spec: StyleSpec) -> dict:
    """The shape the UI renders: palette, fonts, layouts, coverage."""
    present = spec.archetypes_present()
    return {
        "client": spec.meta.client,
        "project": spec.meta.project,
        "spec_version": spec.meta.spec_version,
        "extracted_at": spec.meta.extracted_at,
        "content_slides_ignored": spec.meta.content_slides_ignored,
        "layouts_found": spec.meta.layouts_found,
        "slide_width": spec.slide_width,
        "slide_height": spec.slide_height,
        "colors": spec.theme.colors,
        "fonts": {
            "major_latin": spec.theme.fonts.major_latin,
            "minor_latin": spec.theme.fonts.minor_latin,
            "major_cs": spec.theme.fonts.major_cs,
            "minor_cs": spec.theme.fonts.minor_cs,
        },
        "logo": {
            "present": spec.brand.logo.present,
            "rule": spec.brand.logo.layout_rule,
        },
        "footer": {
            "text": spec.brand.footer.text,
            "slide_numbers": spec.brand.footer.show_slide_number,
        },
        "grid": {
            "margin_left": spec.grid.margin_left_frac,
            "margin_top": spec.grid.margin_top_frac,
            "columns": spec.grid.columns,
            "gutter": spec.grid.gutter_frac,
        },
        "layouts": [
            {"name": l.name, "archetype": l.archetype, "source": l.source,
             "placeholders": len(l.content_placeholders())}
            for l in spec.layouts
        ],
        "archetypes_present": sorted(present),
        "archetypes_missing": sorted(set(archetypes.ALL_ARCHETYPES) - present),
    }


# --- classification preview ----------------------------------------------

@app.post("/api/classify")
async def classify(content: UploadFile = File(...)):
    """
    Read a rough deck and report how each slide would be routed.

    Called before formatting so a designer can see - and correct - the
    routing rather than discovering it in the output.
    """
    path = _save_upload(content, "content.pptx")
    try:
        prs = Presentation(str(path))
        results = classify_deck(prs, prs.slide_width, prs.slide_height)
    except Exception as exc:                                  # noqa: BLE001
        raise HTTPException(422, f"Could not read the deck: {exc}") from exc

    return {
        "slides": [
            {
                "slide": i,
                "archetype": result.archetype,
                "label": archetypes.label_for(result.archetype),
                "confidence": round(result.confidence, 2),
                "evidence": result.evidence,
            }
            for i, result in enumerate(results, start=1)
        ]
    }


# --- jobs -----------------------------------------------------------------

@app.post("/api/jobs")
async def create_job(
    content: UploadFile = File(..., description="The rough content deck"),
    master: UploadFile | None = File(None, description="The designer's master .pptx"),
    client: str | None = Form(None),
    project: str | None = Form(None),
    overrides: str | None = Form(None, description='e.g. {"4": "quote"}'),
    use_banked_spec: bool = Form(False),
):
    """
    Run a formatting job and record it.

    With a master, all three stages run. With `use_banked_spec` and a
    client instead, Stage 1 is skipped and the archived spec is used - the
    repeat-deck path.
    """
    if not master and not use_banked_spec:
        raise HTTPException(
            400, "Provide a master .pptx, or set use_banked_spec with a client."
        )
    if use_banked_spec and not client:
        raise HTTPException(400, "use_banked_spec requires a client.")

    bank, jobs = _bank(), _jobs()
    content_path = _save_upload(content, "content.pptx")
    master_path = _save_upload(master, "master.pptx") if master else None
    parsed_overrides = _parse_overrides(overrides)

    record = jobs.create(
        client=client,
        project=project,
        master_filename=master.filename if master else None,
        content_filename=content.filename,
    )
    output_path = WORK_DIR / f"{record.job_id}_formatted.pptx"

    try:
        if use_banked_spec:
            report = format_with_banked_spec(
                client=client,
                rough_content_pptx_path=content_path,
                output_pptx_path=output_path,
                bank=bank,
                layout_overrides=parsed_overrides,
            )
        else:
            report = run_pipeline(
                master_pptx_path=master_path,
                rough_content_pptx_path=content_path,
                output_pptx_path=output_path,
                bank=bank,
                client=client,
                project=project,
                layout_overrides=parsed_overrides,
            )
    except LookupError as exc:
        jobs.fail(record.job_id, str(exc))
        raise HTTPException(404, str(exc)) from exc
    except Exception as exc:                                  # noqa: BLE001
        jobs.fail(record.job_id, str(exc))
        raise HTTPException(422, f"Formatting failed: {exc}") from exc

    record = jobs.complete(record.job_id, output_path, report)
    return JSONResponse({
        "job": record.summary(),
        "report": report,
        "download_url": f"/api/jobs/{record.job_id}/download",
    })


@app.get("/api/jobs")
def list_jobs(limit: int = 25):
    return {"jobs": [record.summary() for record in _jobs().list(limit=limit)]}


@app.get("/api/jobs/{job_id}")
def get_job(job_id: str):
    record = _jobs().get(job_id)
    if record is None:
        raise HTTPException(404, f"No such job '{job_id}'")
    return {"job": record.summary(), "report": record.report}


@app.get("/api/jobs/{job_id}/download")
def download_job(job_id: str):
    """Download a finished deck. Opens and edits normally in PowerPoint."""
    jobs = _jobs()
    record = jobs.get(job_id)
    if record is None:
        raise HTTPException(404, f"No such job '{job_id}'")
    path = jobs.output_path(job_id)
    if path is None or not Path(path).exists():
        raise HTTPException(404, f"Job '{job_id}' has no output to download")
    return FileResponse(path, filename=record.output_name, media_type=PPTX_MEDIA_TYPE)


# --- template bank --------------------------------------------------------

@app.get("/api/bank")
def bank_list():
    bank = _bank()
    return {
        "entries": [
            {**vars(entry), "has_master": bank.has_master(entry.entry_id)}
            for entry in bank.list_entries()
        ]
    }


@app.get("/api/bank/{entry_id}")
def bank_show(entry_id: str):
    try:
        spec = _bank().load(entry_id)
    except FileNotFoundError as exc:
        raise HTTPException(404, str(exc)) from exc
    return {"summary": _spec_summary(spec), "style_spec": spec.to_dict()}


@app.put("/api/bank/{entry_id}")
def bank_refine(entry_id: str, style_spec: str = Form(...)):
    """
    Fold a designer's corrections back into an archived spec.

    The previous revision is kept, so a correction stays auditable and the
    client's next deck inherits the fix.
    """
    try:
        spec = StyleSpec.from_dict(json.loads(style_spec))
        revision = _bank().refine(entry_id, spec)
    except FileNotFoundError as exc:
        raise HTTPException(404, str(exc)) from exc
    except (json.JSONDecodeError, TypeError) as exc:
        raise HTTPException(400, f"style_spec is not valid Style Spec JSON: {exc}") from exc
    return {"entry_id": entry_id, "revision": revision}
