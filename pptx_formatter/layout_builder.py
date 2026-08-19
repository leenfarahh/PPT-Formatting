"""
Building real slide layouts from Style Spec JSON.

python-pptx has no API for adding a slide layout to a presentation - the
collection is read-only, which is the limitation that used to force this
tool to restyle a generic bundled template instead of extending the
designer's own master. It is only an API gap, though, not a format one:
a layout is just an XML part related to the master and listed in its
`sldLayoutIdLst`. This module writes that part directly.

Two entry points:

*   `add_layout()` takes a LayoutSpec (from the Template Bank, or generated)
    and materializes it as a live layout on a presentation's master.
*   `generated_layout_spec()` synthesizes a LayoutSpec for an archetype from
    the submission's own grid, for the cold-start case where the bank holds
    nothing for that archetype yet.

Geometry arrives as slide fractions and is resolved against the target
deck's dimensions here, which is what lets a layout banked from a 4:3
master land correctly on a 16:9 submission.
"""
from __future__ import annotations

import uuid
from pathlib import Path

from pptx import Presentation
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.parts.slide import SlideLayoutPart

from . import archetypes
from .style_spec import (
    LayoutSpec, PlaceholderSpec, BackgroundSpec, StyleSpec, Grid, SOURCE_GENERATED,
)

A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Conventional placeholder indices for page furniture, matching what
# PowerPoint itself emits.
IDX_DATE, IDX_FOOTER, IDX_SLIDE_NUMBER = 10, 11, 12

# archetype -> OOXML p:sldLayout/@type. PowerPoint uses this to label the
# layout in its own UI, so it's worth setting to the nearest real value.
ARCHETYPE_OOXML_TYPE = {
    archetypes.TITLE_SLIDE: "title",
    archetypes.SECTION_HEADER: "secHead",
    archetypes.TITLE_ONLY: "titleOnly",
    archetypes.TITLE_AND_CONTENT: "obj",
    archetypes.TWO_CONTENT: "twoObj",
    archetypes.THREE_CONTENT: "twoObj",
    archetypes.COMPARISON: "twoTxTwoObj",
    archetypes.QUOTE: "obj",
    archetypes.BIG_STATEMENT: "titleOnly",
    archetypes.PICTURE_FULL: "blank",
    archetypes.PICTURE_CAPTION: "picTx",
    archetypes.TABLE: "tbl",
    archetypes.CHART: "chart",
    archetypes.BLANK: "blank",
    archetypes.CLOSING: "secHead",
}


# --- XML generation -------------------------------------------------------

def _esc(text: str) -> str:
    return (
        (text or "")
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def _attr(name: str, value) -> str:
    """Render an XML attribute, or nothing when the value is None.

    Omission is meaningful here: a placeholder with no `algn` inherits
    alignment from the master, which is what we want unless the spec
    explicitly recorded an override.
    """
    if value is None:
        return ""
    if isinstance(value, bool):
        value = "1" if value else "0"
    return f' {name}="{_esc(str(value))}"'


def _list_style_xml(ph: PlaceholderSpec) -> str:
    """`a:lstStyle` carrying only the properties the spec actually set."""
    p_attrs = _attr("algn", ph.alignment) + _attr("rtl", ph.rtl)
    r_attrs = ""
    if ph.font_size_pt is not None:
        r_attrs += _attr("sz", int(round(ph.font_size_pt * 100)))
    if ph.bold is not None:
        r_attrs += _attr("b", ph.bold)

    if not p_attrs and not r_attrs:
        return "<a:lstStyle/>"
    inner = f"<a:defRPr{r_attrs}/>" if r_attrs else ""
    return f"<a:lstStyle><a:lvl1pPr{p_attrs}>{inner}</a:lvl1pPr></a:lstStyle>"


def _field_run_xml(field_type: str, text: str) -> str:
    """A live PowerPoint field (date or slide number) rather than static text."""
    guid = "{%s}" % str(uuid.uuid4()).upper()
    return (
        f'<a:fld id="{guid}" type="{field_type}">'
        f"<a:t>{_esc(text)}</a:t></a:fld>"
    )


def _placeholder_xml(ph: PlaceholderSpec, shape_id: int, sw: int, sh: int) -> str:
    left, top, width, height = ph.to_emu(sw, sh)

    ph_attrs = _attr("type", ph.ph_type if ph.ph_type != "body" or ph.idx is None else "body")
    # A title placeholder carries no idx; everything else needs one.
    if ph.idx is not None:
        ph_attrs += _attr("idx", ph.idx)

    body_attrs = _attr("anchor", ph.anchor)
    if ph.rtl:
        body_attrs += ' rtlCol="1"'

    # Furniture placeholders carry live fields so PowerPoint keeps them
    # updated; content placeholders get an empty paragraph.
    if ph.ph_type == "sldNum":
        para = f"<a:p><a:pPr{_attr('algn', ph.alignment)}/>{_field_run_xml('slidenum', '‹#›')}</a:p>"
    elif ph.ph_type == "dt":
        para = f"<a:p><a:pPr{_attr('algn', ph.alignment)}/>{_field_run_xml('datetime1', '')}</a:p>"
    else:
        para = '<a:p><a:endParaRPr lang="en-US"/></a:p>'

    return f"""<p:sp>
      <p:nvSpPr>
        <p:cNvPr id="{shape_id}" name="{_esc(ph.name or f'Placeholder {shape_id}')}"/>
        <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
        <p:nvPr><p:ph{ph_attrs}/></p:nvPr>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{left}" y="{top}"/><a:ext cx="{width}" cy="{height}"/></a:xfrm>
      </p:spPr>
      <p:txBody>
        <a:bodyPr{body_attrs}/>
        {_list_style_xml(ph)}
        {para}
      </p:txBody>
    </p:sp>"""


def _background_xml(bg: BackgroundSpec) -> str:
    """
    Inline background XML.

    Image backgrounds are deliberately excluded here: they need a
    relationship to an image part, which can only be added once the part
    exists. `_apply_image_background()` handles those after the fact.
    """
    if bg.kind == "solid" and bg.color_hex:
        fill = f'<a:solidFill><a:srgbClr val="{_esc(bg.color_hex)}"/></a:solidFill>'
    elif bg.kind == "theme" and bg.theme_role:
        fill = f'<a:solidFill><a:schemeClr val="{_esc(bg.theme_role)}"/></a:solidFill>'
    elif bg.kind == "gradient" and bg.gradient_stops:
        stops = bg.gradient_stops
        if len(stops) == 1:
            stops = stops * 2
        gs = "".join(
            f'<a:gs pos="{int(i * 100000 / max(1, len(stops) - 1))}">'
            f'<a:srgbClr val="{_esc(c)}"/></a:gs>'
            for i, c in enumerate(stops)
        )
        fill = (
            f'<a:gradFill rotWithShape="1"><a:gsLst>{gs}</a:gsLst>'
            f'<a:lin ang="5400000" scaled="0"/></a:gradFill>'
        )
    else:
        return ""
    return f"<p:bg><p:bgPr>{fill}<a:effectLst/></p:bgPr></p:bg>"


def build_layout_xml(spec: LayoutSpec, sw: int, sh: int, show_master_shapes: bool = True) -> str:
    """Render a LayoutSpec as a complete `p:sldLayout` document."""
    shapes = []
    for i, ph in enumerate(spec.placeholders):
        shapes.append(_placeholder_xml(ph, shape_id=i + 2, sw=sw, sh=sh))

    show_attr = "" if show_master_shapes else ' showMasterSp="0"'
    ooxml_type = spec.ooxml_type or ARCHETYPE_OOXML_TYPE.get(spec.archetype, "obj")

    return f"""<p:sldLayout {nsdecls("a", "p", "r")} type="{_esc(ooxml_type)}" preserve="1"{show_attr}>
  <p:cSld name="{_esc(spec.name)}">
    {_background_xml(spec.background)}
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="0" y="0"/><a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/>
        </a:xfrm>
      </p:grpSpPr>
      {"".join(shapes)}
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sldLayout>"""


# --- part wiring ----------------------------------------------------------

def _next_layout_partname(package) -> PackURI:
    used = {str(p.partname) for p in package.iter_parts()}
    n = 1
    while f"/ppt/slideLayouts/slideLayout{n}.xml" in used:
        n += 1
    return PackURI(f"/ppt/slideLayouts/slideLayout{n}.xml")


def add_picture_to(shapes, image_path: str, left: int, top: int,
                   width: int, height: int, name: str = "Picture"):
    """
    Add a picture to any shape collection, including a slide master's.

    python-pptx only exposes `add_picture()` on slide shapes; `MasterShapes`
    has no such method. The underlying operation is the same either way -
    relate the image part, then append a `p:pic` to the shape tree - so
    this does it directly, which is what lets the logo go on the master
    where every layout inherits it.
    """
    from pptx.oxml.shapes.picture import CT_Picture

    part = shapes.part
    _, rId = part.get_or_add_image_part(image_path)
    shape_id = max(
        [int(el.get("id")) for el in shapes._spTree.iter(qn("p:cNvPr")) if el.get("id")]
        or [1]
    ) + 1
    pic = CT_Picture.new_pic(shape_id, name, name, rId, left, top, width, height)
    shapes._spTree.append(pic)
    return pic


def _apply_image_background(layout_part, asset_path: str) -> None:
    """Attach a picture-fill background, adding the image part and rel."""
    path = Path(asset_path)
    if not path.exists():
        return
    image_part, rId = layout_part.get_or_add_image_part(str(path))
    cSld = layout_part._element.find(qn("p:cSld"))
    if cSld is None:
        return
    for existing in cSld.findall(qn("p:bg")):
        cSld.remove(existing)
    bg = parse_xml(
        f'<p:bg {nsdecls("a", "p", "r")}><p:bgPr>'
        f'<a:blipFill rotWithShape="1"><a:blip r:embed="{rId}"/>'
        f'<a:stretch><a:fillRect/></a:stretch></a:blipFill>'
        f"<a:effectLst/></p:bgPr></p:bg>"
    )
    cSld.insert(0, bg)


def add_layout(prs: Presentation, spec: LayoutSpec, show_master_shapes: bool = True):
    """
    Materialize a LayoutSpec as a live slide layout on the presentation's
    first master, and return the resulting SlideLayout.

    Creates the layout part, relates it to the master in both directions
    (the layout->master relationship is required by the schema), and
    registers it in the master's `sldLayoutIdLst` so PowerPoint lists it.
    """
    master = prs.slide_masters[0]
    master_part = master.part
    package = prs.part.package

    xml = build_layout_xml(spec, prs.slide_width, prs.slide_height, show_master_shapes)
    layout_part = SlideLayoutPart(
        _next_layout_partname(package), CT.PML_SLIDE_LAYOUT, package, parse_xml(xml)
    )

    layout_part.relate_to(master_part, RT.SLIDE_MASTER)
    rId = master_part.relate_to(layout_part, RT.SLIDE_LAYOUT)

    sldLayoutIdLst = master._element.get_or_add_sldLayoutIdLst()
    entry = sldLayoutIdLst.makeelement(qn("p:sldLayoutId"), {})
    # Layout ids must be >= 2147483648 per the schema.
    existing = [int(e.get("id")) for e in sldLayoutIdLst if e.get("id")]
    entry.set("id", str(max(existing) + 1 if existing else 2147483649))
    entry.set(qn("r:id"), rId)
    sldLayoutIdLst.append(entry)

    if spec.background.kind == "image" and spec.background.asset_path:
        _apply_image_background(layout_part, spec.background.asset_path)

    return layout_part.slide_layout


# --- generated (cold-bank) layouts ---------------------------------------

def _content_box(grid: Grid) -> tuple[float, float, float, float]:
    left = grid.margin_left_frac
    top = grid.margin_top_frac
    width = 1.0 - grid.margin_left_frac - grid.margin_right_frac
    height = 1.0 - grid.margin_top_frac - grid.margin_bottom_frac
    return left, top, width, height


def _furniture(spec: StyleSpec, archetype: str) -> list:
    """Footer and slide-number placeholders, per the brand's field rules."""
    footer_spec = spec.brand.footer
    out = []
    if footer_spec.footer_on(archetype):
        out.append(PlaceholderSpec(
            ph_type="ftr", idx=IDX_FOOTER, name="Footer Placeholder",
            left_frac=footer_spec.left_frac if footer_spec.left_frac is not None else 0.08,
            top_frac=footer_spec.top_frac if footer_spec.top_frac is not None else 0.93,
            width_frac=footer_spec.width_frac if footer_spec.width_frac is not None else 0.5,
            height_frac=footer_spec.height_frac if footer_spec.height_frac is not None else 0.04,
            alignment="l", anchor="ctr",
        ))
    if footer_spec.slide_number_on(archetype):
        out.append(PlaceholderSpec(
            ph_type="sldNum", idx=IDX_SLIDE_NUMBER, name="Slide Number Placeholder",
            left_frac=(footer_spec.slide_number_left_frac
                       if footer_spec.slide_number_left_frac is not None else 0.86),
            top_frac=(footer_spec.slide_number_top_frac
                      if footer_spec.slide_number_top_frac is not None else 0.93),
            width_frac=(footer_spec.slide_number_width_frac
                        if footer_spec.slide_number_width_frac is not None else 0.06),
            height_frac=(footer_spec.slide_number_height_frac
                         if footer_spec.slide_number_height_frac is not None else 0.04),
            alignment="r", anchor="ctr",
        ))
    return out


def generated_layout_spec(archetype: str, spec: StyleSpec) -> LayoutSpec:
    """
    Synthesize a layout for `archetype` from the submission's own grid.

    Used only when the Template Bank has nothing for this archetype - a
    cold bank on the very first submission. The geometry is derived from
    the designer's margins and gutter, so even a generated layout sits on
    their grid rather than on a generic default.
    """
    grid = spec.grid
    left, top, width, height = _content_box(grid)
    gutter = grid.gutter_frac
    title_h = min(0.18, height * 0.25)
    body_top = top + title_h + gutter
    body_h = max(0.1, height - title_h - gutter)
    rtl = None

    def ph(ph_type, idx, name, l, t, w, h, **kw):
        return PlaceholderSpec(
            ph_type=ph_type, idx=idx, name=name,
            left_frac=round(l, 5), top_frac=round(t, 5),
            width_frac=round(w, 5), height_frac=round(h, 5),
            rtl=rtl, **kw
        )

    items: list = []

    if archetype in (archetypes.TITLE_SLIDE, archetypes.CLOSING):
        items = [
            ph("ctrTitle", None, "Title", left, 0.34, width, 0.22, anchor="b"),
            ph("subTitle", 1, "Subtitle", left, 0.58, width, 0.14, anchor="t"),
        ]
    elif archetype == archetypes.SECTION_HEADER:
        items = [
            ph("title", None, "Title", left, 0.38, width, 0.18, anchor="b"),
            ph("body", 1, "Text Placeholder", left, 0.58, width, 0.12, anchor="t"),
        ]
    elif archetype == archetypes.TITLE_ONLY:
        items = [ph("title", None, "Title", left, top, width, title_h, anchor="b")]
    elif archetype == archetypes.BIG_STATEMENT:
        items = [ph("title", None, "Statement", left, 0.3, width, 0.4,
                    anchor="ctr", alignment="ctr")]
    elif archetype == archetypes.QUOTE:
        items = [
            ph("title", None, "Quote", left, 0.28, width, 0.32,
               anchor="ctr", alignment="ctr"),
            ph("body", 1, "Attribution", left, 0.63, width, 0.1,
               anchor="t", alignment="ctr"),
        ]
    elif archetype == archetypes.TITLE_AND_CONTENT:
        items = [
            ph("title", None, "Title", left, top, width, title_h, anchor="b"),
            ph("body", 1, "Content Placeholder", left, body_top, width, body_h, anchor="t"),
        ]
    elif archetype in (archetypes.TWO_CONTENT, archetypes.THREE_CONTENT):
        n = 2 if archetype == archetypes.TWO_CONTENT else 3
        col_w = (width - gutter * (n - 1)) / n
        items = [ph("title", None, "Title", left, top, width, title_h, anchor="b")]
        for i in range(n):
            items.append(ph("body", i + 1, f"Content Placeholder {i + 1}",
                            left + i * (col_w + gutter), body_top, col_w, body_h, anchor="t"))
    elif archetype == archetypes.COMPARISON:
        col_w = (width - gutter) / 2
        head_h = min(0.08, body_h * 0.22)
        items = [ph("title", None, "Title", left, top, width, title_h, anchor="b")]
        for i in range(2):
            x = left + i * (col_w + gutter)
            items.append(ph("body", i * 2 + 1, f"Heading {i + 1}",
                            x, body_top, col_w, head_h, anchor="b", bold=True))
            items.append(ph("body", i * 2 + 2, f"Content {i + 1}",
                            x, body_top + head_h + gutter / 2, col_w,
                            body_h - head_h - gutter / 2, anchor="t"))
    elif archetype == archetypes.PICTURE_FULL:
        items = [ph("pic", 1, "Picture Placeholder", 0.0, 0.0, 1.0, 1.0)]
    elif archetype == archetypes.PICTURE_CAPTION:
        pic_w = width * 0.55
        text_w = width - pic_w - gutter
        items = [
            ph("pic", 1, "Picture Placeholder", left, top, pic_w, height),
            ph("title", None, "Title", left + pic_w + gutter, top, text_w, title_h, anchor="b"),
            ph("body", 2, "Caption", left + pic_w + gutter, body_top, text_w, body_h, anchor="t"),
        ]
    elif archetype == archetypes.TABLE:
        items = [
            ph("title", None, "Title", left, top, width, title_h, anchor="b"),
            ph("tbl", 1, "Table Placeholder", left, body_top, width, body_h),
        ]
    elif archetype == archetypes.CHART:
        items = [
            ph("title", None, "Title", left, top, width, title_h, anchor="b"),
            ph("chart", 1, "Chart Placeholder", left, body_top, width, body_h),
        ]
    elif archetype == archetypes.BLANK:
        items = []
    else:
        items = [
            ph("title", None, "Title", left, top, width, title_h, anchor="b"),
            ph("body", 1, "Content Placeholder", left, body_top, width, body_h, anchor="t"),
        ]

    items.extend(_furniture(spec, archetype))

    return LayoutSpec(
        name=archetypes.label_for(archetype),
        archetype=archetype,
        ooxml_type=ARCHETYPE_OOXML_TYPE.get(archetype, "obj"),
        placeholders=items,
        background=BackgroundSpec(kind="inherit"),
        source=SOURCE_GENERATED,
    )
