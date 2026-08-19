"""
Validation checks run over the finished deck.

These are static checks read straight from the object model - no rendering
step - so they're cheap enough to run on every slide. That also bounds what
they can know: text overflow in particular is estimated from character
counts, because measuring it properly requires laying out the text in the
actual font, which needs a rendering pass.

Findings are advisory. They surface in the pipeline report for a designer
to review; nothing here modifies the deck.
"""
from __future__ import annotations

from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Emu

from . import rtl

# WCAG AA for normal-size text.
MIN_CONTRAST = 4.5
# Characters per line at a nominal 18pt in a full-width box; scaled by the
# shape's actual width and font size in the overflow estimate.
BASE_CHARS_PER_LINE = 90


def _relative_luminance(rgb: tuple) -> float:
    """WCAG relative luminance. https://www.w3.org/TR/WCAG21/#dfn-relative-luminance"""
    def channel(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = (channel(c) for c in rgb)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def contrast_ratio(hex_a: str, hex_b: str) -> float:
    def to_rgb(h):
        h = h.lstrip("#")
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    l1 = _relative_luminance(to_rgb(hex_a)) + 0.05
    l2 = _relative_luminance(to_rgb(hex_b)) + 0.05
    return max(l1, l2) / min(l1, l2)


def check_contrast(slide, background_hex: str, min_ratio: float = MIN_CONTRAST) -> list:
    """
    Flag text whose literal color fails WCAG AA against the background.

    Only runs with an explicit RGB are checked. Text using a theme
    reference is assumed fine, since the theme is what Stage 2 controls and
    checking it here would just re-report the brand's own palette.
    """
    issues = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                try:
                    hex_color = str(run.font.color.rgb)
                except (AttributeError, TypeError):
                    continue
                ratio = contrast_ratio(hex_color, background_hex)
                if ratio < min_ratio:
                    issues.append(
                        f"Low contrast on '{shape.name}': #{hex_color} on #{background_hex} "
                        f"= {ratio:.2f}:1 (needs {min_ratio}:1)"
                    )
    return issues


def check_overflow(slide, slide_width: int) -> list:
    """
    Estimate whether text exceeds its shape.

    Deliberately conservative: line count is estimated from character
    counts scaled by the shape's width and font size, so it will miss some
    real overflows and occasionally flag text that fits. A precise answer
    needs a real text-layout pass in the rendering font.

    Footers, dates and slide numbers are skipped. Their size is set by the
    master and never written on the run, so the estimate would have to
    guess a size for them - and guessing body size for a footer flags every
    slide in the deck, which buries the findings that matter.
    """
    issues = []
    latent = {PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER}
    for shape in slide.shapes:
        if not shape.has_text_frame or not shape.height or not shape.width:
            continue
        if shape.is_placeholder and shape.placeholder_format.type in latent:
            continue
        text = shape.text_frame.text
        if not text.strip():
            continue

        sizes = [
            run.font.size.pt
            for para in shape.text_frame.paragraphs
            for run in para.runs
            if run.font.size is not None
        ]
        font_pt = max(sizes) if sizes else 18.0
        width_ratio = shape.width / slide_width if slide_width else 1.0
        chars_per_line = max(8, int(BASE_CHARS_PER_LINE * width_ratio * (18.0 / font_pt)))

        lines = 0
        for line in text.split("\n"):
            lines += max(1, -(-len(line) // chars_per_line))

        line_height = Emu(int(font_pt * 12700 * 1.2))     # 12700 EMU per point
        estimated = lines * line_height
        if estimated > shape.height * 1.15:               # 15% slack
            issues.append(
                f"'{shape.name}' may overflow: ~{lines} lines at {font_pt:.0f}pt needs "
                f"~{Emu(int(estimated)).inches:.2f}in, shape is "
                f"{Emu(int(shape.height)).inches:.2f}in tall"
            )
    return issues


def check_off_slide(slide, slide_width: int, slide_height: int) -> list:
    """Flag shapes extending past the slide edges."""
    issues = []
    for shape in slide.shapes:
        if shape.left is None or shape.top is None:
            continue
        right = shape.left + (shape.width or 0)
        bottom = shape.top + (shape.height or 0)
        if shape.left < 0 or shape.top < 0 or right > slide_width or bottom > slide_height:
            issues.append(f"'{shape.name}' extends beyond the slide boundary")
    return issues


def check_empty_placeholders(slide) -> list:
    """
    Flag placeholders left empty.

    An empty placeholder prints as nothing but shows prompt text while
    editing, so it reads as unfinished to whoever opens the deck.
    """
    issues = []
    latent = {PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER}
    for ph in slide.placeholders:
        if ph.placeholder_format.type in latent:
            continue
        if ph.has_text_frame and not ph.text_frame.text.strip():
            issues.append(f"'{ph.name}' is an empty placeholder")
    return issues


def check_mixed_direction(slide) -> list:
    """
    Flag frames mixing right-to-left and left-to-right paragraphs.

    Legitimate in a bilingual deck, but worth a designer's eye: mixed
    direction inside one frame is also what a bad copy-paste looks like.
    """
    issues = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        directions = set()
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs)
            if text.strip():
                directions.add(rtl.is_rtl_text(text))
        if len(directions) > 1:
            issues.append(f"'{shape.name}' mixes right-to-left and left-to-right paragraphs")
    return issues


def check_slide(slide, spec, background_hex: str) -> list:
    """Run every check over one slide."""
    issues = []
    issues.extend(check_contrast(slide, background_hex))
    issues.extend(check_overflow(slide, spec.slide_width))
    issues.extend(check_off_slide(slide, spec.slide_width, spec.slide_height))
    issues.extend(check_empty_placeholders(slide))
    issues.extend(check_mixed_direction(slide))
    return issues
