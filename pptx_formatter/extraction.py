"""
Stage 1 - Master ingestion and style extraction.

Reads a designer-submitted master `.pptx` and produces a complete Style
Spec: theme colors, per-script typefaces, backgrounds, the logo (bytes and
placement rule), footer and page-number rules, a derived grid, chart/table
styling, an icon palette, and a LayoutSpec for every layout the designer
authored.

Two things worth calling out:

*   **Content slides in the submission are ignored.** The designer may hand
    over an empty master or a fully populated deck; either way only
    `slideMasters` and `slideLayouts` are read. The count of skipped slides
    is recorded in `meta.content_slides_ignored` so the caller can report it.

*   **Placeholder geometry is resolved through inheritance.** A layout
    placeholder usually carries no `<a:xfrm>` of its own and inherits
    position from the master placeholder of the same type. python-pptx's
    `LayoutPlaceholder` resolves that, so the fractions banked here are the
    placeholder's *effective* geometry rather than an empty override.

python-pptx exposes no public API for theme colors/fonts, backgrounds, or
header/footer rules, so those are read from the XML parts via lxml.
"""
from __future__ import annotations

import re
from datetime import datetime, timezone
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.enum.shapes import MSO_SHAPE_TYPE

from . import archetypes
from .style_spec import (
    SPEC_VERSION, StyleSpec, Theme, ThemeFonts, Brand, BrandLogo, BrandFooter,
    BackgroundSpec, LayoutSpec, PlaceholderSpec, Grid, ChartStyle, TableStyle,
    Meta, THEME_COLOR_ROLES, ACCENT_ROLES, SOURCE_DESIGNER,
    LOGO_ALL, LOGO_EXCEPT_TITLE, LOGO_NONE,
)

A_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
P_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

# Placeholder types that are page furniture rather than content.
FURNITURE = {"dt", "ftr", "sldNum"}


# --- theme ---------------------------------------------------------------

def get_theme_part(prs: Presentation):
    """The OPC Part holding the (first) master's theme."""
    return prs.slide_masters[0].part.part_related_by(RT.THEME)


def get_theme_root(prs: Presentation) -> etree._Element:
    return etree.fromstring(get_theme_part(prs).blob)


def _extract_theme_colors(theme_root: etree._Element) -> dict:
    colors = {}
    scheme = theme_root.find(".//a:clrScheme", A_NS)
    if scheme is None:
        return colors
    for role in THEME_COLOR_ROLES:
        role_el = scheme.find(f"a:{role}", A_NS)
        if role_el is None:
            continue
        srgb = role_el.find("a:srgbClr", A_NS)
        if srgb is not None and srgb.get("val"):
            colors[role] = srgb.get("val").upper()
            continue
        sys_clr = role_el.find("a:sysClr", A_NS)
        if sys_clr is not None and sys_clr.get("lastClr"):
            colors[role] = sys_clr.get("lastClr").upper()
    return colors


def _extract_theme_fonts(theme_root: etree._Element) -> ThemeFonts:
    """Pull the Latin, East Asian and complex-script (Arabic) typefaces."""
    fonts = ThemeFonts()
    scheme = theme_root.find(".//a:fontScheme", A_NS)
    if scheme is None:
        return fonts
    for prefix in ("major", "minor"):
        for script, attr in (("latin", "latin"), ("ea", "ea"), ("cs", "cs")):
            el = scheme.find(f"a:{prefix}Font/a:{script}", A_NS)
            if el is not None and el.get("typeface"):
                setattr(fonts, f"{prefix}_{attr}", el.get("typeface"))
    return fonts


# --- backgrounds ---------------------------------------------------------

def _extract_background(cSld_el, part, asset_dir: Path | None, tag: str) -> BackgroundSpec:
    """
    Read a `p:cSld/p:bg` into a BackgroundSpec.

    Absence of `p:bg` means "inherit", the common and correct case for a
    layout sitting under a styled master - recording that keeps us from
    baking the master's background into every layout.
    """
    if cSld_el is None:
        return BackgroundSpec(kind="inherit")

    bg = cSld_el.find("p:bg", P_NS)
    if bg is None:
        return BackgroundSpec(kind="inherit")

    # A background is either a direct fill (bgPr) or a reference into the
    # theme's format scheme (bgRef). A bgRef follows the theme, which Stage 2
    # rewrites anyway, so it stays "inherit".
    bgPr = bg.find("p:bgPr", P_NS)
    if bgPr is None:
        return BackgroundSpec(kind="inherit")

    solid = bgPr.find("a:solidFill", A_NS)
    if solid is not None:
        srgb = solid.find("a:srgbClr", A_NS)
        if srgb is not None and srgb.get("val"):
            return BackgroundSpec(kind="solid", color_hex=srgb.get("val").upper())
        scheme_clr = solid.find("a:schemeClr", A_NS)
        if scheme_clr is not None and scheme_clr.get("val"):
            return BackgroundSpec(kind="theme", theme_role=scheme_clr.get("val"))

    grad = bgPr.find("a:gradFill", A_NS)
    if grad is not None:
        stops = []
        for gs in grad.findall(".//a:gs", A_NS):
            srgb = gs.find("a:srgbClr", A_NS)
            if srgb is not None and srgb.get("val"):
                stops.append(srgb.get("val").upper())
        return BackgroundSpec(kind="gradient", gradient_stops=stops,
                              color_hex=stops[0] if stops else None)

    blip_fill = bgPr.find("a:blipFill", A_NS)
    if blip_fill is not None and asset_dir is not None:
        blip = blip_fill.find(".//a:blip", A_NS)
        embed = blip.get(f"{{{P_NS['r']}}}embed") if blip is not None else None
        if embed:
            try:
                image_part = part.related_part(embed)
                path = _write_asset(asset_dir, f"{tag}_bg", image_part.blob,
                                    getattr(image_part, "ext", "png"))
                return BackgroundSpec(kind="image", asset_path=str(path))
            except KeyError:
                pass

    return BackgroundSpec(kind="inherit")


def _write_asset(asset_dir: Path, stem: str, blob: bytes, ext: str) -> Path:
    asset_dir.mkdir(parents=True, exist_ok=True)
    ext = (ext or "png").lstrip(".")
    path = asset_dir / f"{stem}.{ext}"
    path.write_bytes(blob)
    return path


# --- logo & footer -------------------------------------------------------

def _find_logo(prs: Presentation, asset_dir: Path | None) -> BrandLogo:
    """
    Pick the logo from the pictures sitting on the slide master, and work
    out which layouts it appears on.

    Detection heuristic: a logo is small and pushed toward an edge. Each
    picture is scored on area (smaller is better) and distance from the
    nearest corner (closer is better). A picture covering a quarter of the
    slide or more is treated as artwork, not a logo.

    Placement rule: a layout with `showMasterSp="0"` suppresses master
    shapes, so the logo does not appear there. Reading that attribute gives
    us the designer's actual intent instead of a guess.
    """
    master = prs.slide_masters[0]
    sw, sh = prs.slide_width, prs.slide_height
    best, best_score = None, float("inf")

    for shape in master.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if None in (shape.left, shape.top, shape.width, shape.height):
            continue
        area_frac = (shape.width * shape.height) / float(sw * sh)
        if area_frac > 0.25:
            continue     # full-bleed artwork, not a logo
        cx = (shape.left + shape.width / 2) / sw
        cy = (shape.top + shape.height / 2) / sh
        corner_dist = min(
            (cx ** 2 + cy ** 2) ** 0.5,
            ((1 - cx) ** 2 + cy ** 2) ** 0.5,
            (cx ** 2 + (1 - cy) ** 2) ** 0.5,
            ((1 - cx) ** 2 + (1 - cy) ** 2) ** 0.5,
        )
        score = area_frac * 4 + corner_dist
        if score < best_score:
            best, best_score = shape, score

    if best is None:
        return BrandLogo(layout_rule=LOGO_NONE)

    asset_path = None
    if asset_dir is not None:
        try:
            image = best.image
            asset_path = str(_write_asset(asset_dir, "logo", image.blob, image.ext))
        except (AttributeError, ValueError):
            asset_path = None

    return BrandLogo(
        asset_path=asset_path,
        left_frac=best.left / sw,
        top_frac=best.top / sh,
        width_frac=best.width / sw,
        height_frac=best.height / sh,
        layout_rule=LOGO_ALL,        # refined once layouts are classified
        exclude_archetypes=[],
    )


def _refine_logo_rule(logo: BrandLogo, prs: Presentation, layouts: list) -> None:
    """
    Turn per-layout `showMasterSp` flags into a logo placement rule.

    Done after layout classification because the rule is expressed in
    archetypes, not layout names.
    """
    if not logo.present:
        return
    suppressed = set()
    layout_els = [l for m in prs.slide_masters for l in m.slide_layouts]
    for layout_el, spec in zip(layout_els, layouts):
        if layout_el._element.get("showMasterSp") == "0":
            suppressed.add(spec.archetype)

    if not suppressed:
        logo.layout_rule = LOGO_ALL
        return
    cover_like = {archetypes.TITLE_SLIDE, archetypes.CLOSING, archetypes.PICTURE_FULL}
    if suppressed <= cover_like:
        logo.layout_rule = LOGO_EXCEPT_TITLE
    else:
        logo.layout_rule = LOGO_ALL
        logo.exclude_archetypes = sorted(suppressed)


def _find_footer(prs: Presentation) -> BrandFooter:
    """
    Read footer text, page-number behavior and date settings.

    The master's `<p:hf>` element carries the designer's show/hide intent
    for each field; its attributes default to true when absent, matching
    the OOXML spec.
    """
    master = prs.slide_masters[0]
    sw, sh = prs.slide_width, prs.slide_height
    footer = BrandFooter()

    hf = master._element.find("p:hf", P_NS)
    if hf is not None:
        def flag(attr, default=True):
            raw = hf.get(attr)
            return default if raw is None else raw in ("1", "true")
        footer.show_footer = flag("ftr")
        footer.show_slide_number = flag("sldNum")
        footer.show_date = flag("dt", default=False)

    for ph in master.placeholders:
        ph_type = str(ph.placeholder_format.type or "")
        geom = None
        if None not in (ph.left, ph.top, ph.width, ph.height):
            geom = (ph.left / sw, ph.top / sh, ph.width / sw, ph.height / sh)

        if "FOOTER" in ph_type:
            if ph.has_text_frame and ph.text_frame.text.strip():
                footer.text = ph.text_frame.text
            if geom:
                (footer.left_frac, footer.top_frac,
                 footer.width_frac, footer.height_frac) = geom
        elif "SLIDE_NUMBER" in ph_type:
            if geom:
                (footer.slide_number_left_frac, footer.slide_number_top_frac,
                 footer.slide_number_width_frac, footer.slide_number_height_frac) = geom
        elif "DATE" in ph_type:
            fld = ph._element.find(".//a:fld", A_NS) if ph.has_text_frame else None
            if fld is not None and fld.get("type"):
                footer.date_format = fld.get("type")

    return footer


# --- layouts -------------------------------------------------------------

def _ph_element(shape):
    """The `p:ph` element for a placeholder shape, or None."""
    return shape._element.find(".//p:nvSpPr/p:nvPr/p:ph", P_NS)


def _read_text_props(shape) -> dict:
    """
    Read typography actually declared on this placeholder.

    Anything absent stays None, meaning "inherits from the master". That
    distinction is load-bearing: the inheritance check can only tell a real
    override from an inherited value if we recorded which were explicit.
    """
    props = {"font_size_pt": None, "bold": None, "alignment": None,
             "anchor": None, "rtl": None}
    if not shape.has_text_frame:
        return props

    tx_body = shape._element.find(".//p:txBody", P_NS)
    if tx_body is None:
        return props

    body_pr = tx_body.find("a:bodyPr", A_NS)
    if body_pr is not None and body_pr.get("anchor"):
        props["anchor"] = body_pr.get("anchor")

    # Prefer the layout's own list style (lvl1pPr); fall back to the first
    # paragraph's properties.
    p_pr = tx_body.find("a:lstStyle/a:lvl1pPr", A_NS)
    def_rpr = None
    if p_pr is not None:
        def_rpr = p_pr.find("a:defRPr", A_NS)
    else:
        first_p = tx_body.find("a:p", A_NS)
        if first_p is not None:
            p_pr = first_p.find("a:pPr", A_NS)
            if p_pr is not None:
                def_rpr = p_pr.find("a:defRPr", A_NS)

    if p_pr is not None:
        if p_pr.get("algn"):
            props["alignment"] = p_pr.get("algn")
        if p_pr.get("rtl") is not None:
            props["rtl"] = p_pr.get("rtl") in ("1", "true")

    if def_rpr is not None:
        if def_rpr.get("sz"):
            props["font_size_pt"] = int(def_rpr.get("sz")) / 100.0
        if def_rpr.get("b") is not None:
            props["bold"] = def_rpr.get("b") in ("1", "true")

    return props


def _extract_placeholders(layout, sw: int, sh: int) -> list:
    specs = []
    for shape in layout.placeholders:
        ph_el = _ph_element(shape)
        ph_type = (ph_el.get("type") if ph_el is not None else None) or "body"
        idx_raw = ph_el.get("idx") if ph_el is not None else None
        idx = int(idx_raw) if idx_raw is not None else None

        # Geometry resolves through master inheritance (LayoutPlaceholder);
        # if even the master has nothing, skip rather than bank zeros.
        if None in (shape.left, shape.top, shape.width, shape.height):
            continue

        specs.append(PlaceholderSpec(
            ph_type=ph_type,
            idx=idx,
            name=shape.name,
            left_frac=shape.left / sw,
            top_frac=shape.top / sh,
            width_frac=shape.width / sw,
            height_frac=shape.height / sh,
            **_read_text_props(shape),
        ))
    return specs


# Names designers actually use, mapped to archetypes. Checked before the
# structural fallback: a designer naming a layout "Quote" is much stronger
# evidence than its placeholder count.
_NAME_PATTERNS = [
    (r"\b(cover|title\s*slide|opening)\b", archetypes.TITLE_SLIDE),
    (r"\b(closing|thank\s*you|end\s*slide|back\s*cover)\b", archetypes.CLOSING),
    (r"\b(section|divider|chapter|break)\b", archetypes.SECTION_HEADER),
    (r"\b(quote|testimonial|pull\s*quote)\b", archetypes.QUOTE),
    (r"\b(stat|statistic|big\s*number|kpi|metric|headline)\b", archetypes.BIG_STATEMENT),
    (r"\b(full\s*bleed|full\s*image|image\s*only|photo\s*full)\b", archetypes.PICTURE_FULL),
    (r"\b(picture|image|photo)\b.*\b(caption|text)\b", archetypes.PICTURE_CAPTION),
    (r"\b(table|matrix)\b", archetypes.TABLE),
    (r"\b(chart|graph|data)\b", archetypes.CHART),
    (r"\b(compar|versus|vs\.?)\b", archetypes.COMPARISON),
    (r"\b(three|3)\s*(column|content|up)\b", archetypes.THREE_CONTENT),
    (r"\b(two|2)\s*(column|content|up)\b", archetypes.TWO_CONTENT),
    (r"\b(agenda|contents|toc)\b", archetypes.TITLE_AND_CONTENT),
    (r"\bblank\b", archetypes.BLANK),
    (r"\btitle\s*only\b", archetypes.TITLE_ONLY),
]

# OOXML p:sldLayout/@type -> archetype.
_OOXML_TYPE_MAP = {
    "title": archetypes.TITLE_SLIDE,
    "secHead": archetypes.SECTION_HEADER,
    "titleOnly": archetypes.TITLE_ONLY,
    "obj": archetypes.TITLE_AND_CONTENT,
    "objOnly": archetypes.TITLE_AND_CONTENT,
    "tx": archetypes.TITLE_AND_CONTENT,
    "twoObj": archetypes.TWO_CONTENT,
    "twoTxTwoObj": archetypes.COMPARISON,
    "twoObjAndTx": archetypes.COMPARISON,
    "fourObj": archetypes.COMPARISON,
    "objTx": archetypes.PICTURE_CAPTION,
    "picTx": archetypes.PICTURE_CAPTION,
    "tbl": archetypes.TABLE,
    "chart": archetypes.CHART,
    "blank": archetypes.BLANK,
    "vertTx": archetypes.TITLE_AND_CONTENT,
    "vertTitleAndTx": archetypes.TITLE_AND_CONTENT,
}


def classify_layout(name: str, ooxml_type: str, placeholders: list) -> str:
    """
    Tag a layout with an archetype, in descending order of evidence
    strength: the designer's own layout name, then a decisive geometric
    signal, then the OOXML layout type, then placeholder composition.
    """
    lowered = (name or "").lower()
    for pattern, archetype in _NAME_PATTERNS:
        if re.search(pattern, lowered):
            return archetype

    content = [p for p in placeholders if p.ph_type not in FURNITURE]
    bodies = [p for p in content if p.ph_type in ("body", "obj", "subTitle")]
    pics = [p for p in content if p.ph_type == "pic"]
    titles = [p for p in content if p.ph_type in ("title", "ctrTitle")]

    # A near-full-slide picture is a full-bleed layout whatever it's called.
    for p in pics:
        if p.width_frac > 0.9 and p.height_frac > 0.9:
            return archetypes.PICTURE_FULL

    if ooxml_type in _OOXML_TYPE_MAP:
        mapped = _OOXML_TYPE_MAP[ooxml_type]
        # Trust the type unless the placeholder count clearly contradicts it
        # (a "twoObj" layout carrying three bodies is a three-column layout
        # the designer repurposed).
        if mapped == archetypes.TWO_CONTENT and len(bodies) >= 3:
            return archetypes.THREE_CONTENT
        return mapped

    if not content:
        return archetypes.BLANK
    if pics and len(bodies) <= 1:
        return archetypes.PICTURE_CAPTION
    if len(bodies) >= 3:
        return archetypes.THREE_CONTENT
    if len(bodies) == 2:
        return archetypes.TWO_CONTENT
    if titles and not bodies:
        return archetypes.TITLE_ONLY
    return archetypes.TITLE_AND_CONTENT


def _extract_layouts(prs: Presentation, asset_dir: Path | None) -> list:
    sw, sh = prs.slide_width, prs.slide_height
    specs = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            ooxml_type = layout._element.get("type") or "obj"
            placeholders = _extract_placeholders(layout, sw, sh)
            specs.append(LayoutSpec(
                name=layout.name,
                archetype=classify_layout(layout.name, ooxml_type, placeholders),
                ooxml_type=ooxml_type,
                placeholders=placeholders,
                background=_extract_background(
                    layout._element.find("p:cSld", P_NS), layout.part, asset_dir,
                    f"layout_{_slug(layout.name)}",
                ),
                source=SOURCE_DESIGNER,
            ))
    return specs


def _slug(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", "_", (text or "layout").lower()).strip("_") or "layout"


# --- grid ----------------------------------------------------------------

# Fewest content placeholders that make an inferred grid worth trusting.
# Three is the point at which a left edge, a right edge and a gutter can all
# be corroborated by more than one layout.
MIN_GRID_SAMPLE = 3

# No inferred margin may exceed this. A designer can of course reserve a
# fifth of the slide, but a margin that wide inferred from placeholder
# extents is far more often an artefact than an intention.
MAX_MARGIN_FRAC = 0.12

# How much wider than the left margin the right one may be inferred before
# it is treated as an artefact and mirrored from the left.
#
# The right margin is derived from the widest placeholder's right edge, so it
# is only meaningful when the sample actually contains a full-width element.
# When it doesn't, the "margin" is really just the gap beyond the widest
# thing that happened to be measured, and it strands a band of dead slide
# down the side of every generated layout. Horizontal grids are near
# symmetric almost always, so a large asymmetry is evidence of a gap in the
# sample rather than of the designer's intent. Vertical margins are left
# alone: a tall title zone over a slim footer band is genuinely common.
MARGIN_ASYMMETRY_LIMIT = 2.0


def derive_grid(layouts: list) -> Grid:
    """
    Infer margins, gutters and guide positions from where the designer
    actually put their placeholders, rather than assuming a default.

    Left/top margins come from the smallest placeholder offset across all
    layouts (the designer's true content edge); right/bottom from the
    largest extent. The column count is inferred from the narrowest content
    placeholder - a half-width body implies a 2-up grid, a third implies
    3-up - then expressed on a 12-column grid, which divides evenly by 2, 3,
    4 and 6 so all those rhythms land on real guide lines.
    """
    grid = Grid()
    content = [
        p
        for layout in layouts
        for p in layout.placeholders
        if p.ph_type not in FURNITURE and p.width_frac > 0.05
    ]
    # Inference needs a representative sample. A master that carries its
    # design in plain shapes rather than placeholders offers a handful of
    # stragglers, and margins derived from those describe nothing: they
    # land on the clamp below and become a content box unrelated to the
    # design, which every generated layout then inherits. Below the
    # threshold the defaults are the more honest answer.
    if len(content) < MIN_GRID_SAMPLE:
        grid.compute_guides()
        return grid

    grid.margin_left_frac = round(min(p.left_frac for p in content), 4)
    grid.margin_top_frac = round(min(p.top_frac for p in content), 4)
    grid.margin_right_frac = round(1.0 - max(p.left_frac + p.width_frac for p in content), 4)
    grid.margin_bottom_frac = round(1.0 - max(p.top_frac + p.height_frac for p in content), 4)

    # No full-width element in the sample means no evidence for the right
    # margin. Mirror the left, which at least keeps the content box centred
    # instead of hard against one edge.
    if (grid.margin_left_frac > 0.01
            and grid.margin_right_frac > grid.margin_left_frac * MARGIN_ASYMMETRY_LIMIT):
        grid.margin_right_frac = grid.margin_left_frac

    # Clamp both ends. A full-bleed layout drags a margin to zero and
    # defeats grid snapping for the whole deck; one stray narrow placeholder
    # drags the opposite margin wide enough to strand a band of dead slide.
    for attr in ("margin_left_frac", "margin_top_frac",
                 "margin_right_frac", "margin_bottom_frac"):
        setattr(grid, attr, max(0.0, min(getattr(grid, attr), MAX_MARGIN_FRAC)))

    # Gutter: the smallest horizontal gap between side-by-side placeholders.
    gutters = []
    for layout in layouts:
        row = sorted(
            (p for p in layout.placeholders if p.ph_type not in FURNITURE),
            key=lambda p: p.left_frac,
        )
        for a, b in zip(row, row[1:]):
            gap = b.left_frac - (a.left_frac + a.width_frac)
            # Only count real side-by-side pairs that vertically overlap.
            if gap > 0.002 and a.top_frac < b.top_frac + b.height_frac and \
               b.top_frac < a.top_frac + a.height_frac:
                gutters.append(gap)
    if gutters:
        grid.gutter_frac = round(min(gutters), 4)

    # Vertical rhythm: the smallest gap between stacked placeholder tops.
    tops = sorted({round(p.top_frac, 4) for p in content})
    gaps = [b - a for a, b in zip(tops, tops[1:]) if b - a > 0.01]
    usable_h = 1.0 - grid.margin_top_frac - grid.margin_bottom_frac
    if gaps and usable_h > 0:
        grid.rows = max(4, min(24, int(round(usable_h / min(gaps)))))

    grid.compute_guides()
    return grid


# --- chart / table / icon styling ---------------------------------------

def derive_chart_style(colors: dict, fonts: ThemeFonts) -> ChartStyle:
    """
    Chart styling implied by the brand's theme.

    Series rotate through the accents in theme order (the convention
    PowerPoint itself uses), gridlines take a light tint of the secondary
    dark color, and axis labels use the body face.
    """
    return ChartStyle(
        series_colors=[colors[r] for r in ACCENT_ROLES if r in colors],
        font=fonts.minor_latin,
        gridline_color=colors.get("lt2") or colors.get("dk2"),
        axis_color=colors.get("dk2") or colors.get("dk1"),
    )


def derive_table_style(colors: dict, fonts: ThemeFonts) -> TableStyle:
    """Table styling implied by the brand's theme."""
    header_fill = colors.get("accent1")
    return TableStyle(
        header_fill=header_fill,
        header_font_color=_readable_on(header_fill, colors),
        body_font_color=colors.get("dk1"),
        banded_fill=colors.get("lt2"),
        border_color=colors.get("lt2") or colors.get("dk2"),
        font=fonts.minor_latin,
    )


def derive_icon_palette(colors: dict) -> list:
    """
    Accent colors approved for recoloring monochrome icons.

    Only the accents: recoloring an icon to the background role would make
    it invisible, and to the body-text role would make it read as text.
    """
    return [colors[r] for r in ACCENT_ROLES if r in colors]


def _readable_on(bg_hex: str | None, colors: dict) -> str | None:
    """Pick the light or dark theme text color that reads on `bg_hex`."""
    if not bg_hex:
        return colors.get("lt1")
    r, g, b = (int(bg_hex[i:i + 2], 16) for i in (0, 2, 4))
    # Perceived luminance (ITU-R BT.601), enough to choose between two options.
    luma = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    return colors.get("dk1", "000000") if luma > 0.6 else colors.get("lt1", "FFFFFF")


# --- entry point ---------------------------------------------------------

def count_ignored_slides(master_pptx_path: str | Path) -> int:
    """How many content slides the submitted master carried (all ignored)."""
    return len(Presentation(str(master_pptx_path)).slides._sldIdLst)


def extract_style_spec(
    master_pptx_path: str | Path,
    asset_dir: str | Path | None = None,
    client: str | None = None,
    project: str | None = None,
) -> StyleSpec:
    """
    Stage 1 entry point: a master `.pptx` becomes a Style Spec.

    `asset_dir` is where extracted binaries (the logo, any background
    images) are written. Without it the spec still records the logo's
    geometry but carries no bytes, so Stage 2 cannot re-insert it.
    """
    master_pptx_path = Path(master_pptx_path)
    asset_dir = Path(asset_dir) if asset_dir else None

    prs = Presentation(str(master_pptx_path))
    theme_root = get_theme_root(prs)
    master = prs.slide_masters[0]

    colors = _extract_theme_colors(theme_root)
    fonts = _extract_theme_fonts(theme_root)
    layouts = _extract_layouts(prs, asset_dir)

    logo = _find_logo(prs, asset_dir)
    _refine_logo_rule(logo, prs, layouts)

    spec = StyleSpec(
        theme=Theme(colors=colors, fonts=fonts),
        brand=Brand(logo=logo, footer=_find_footer(prs)),
        layouts=layouts,
        master_background=_extract_background(
            master._element.find("p:cSld", P_NS), master.part, asset_dir, "master",
        ),
        chart_style=derive_chart_style(colors, fonts),
        table_style=derive_table_style(colors, fonts),
        icon_palette=derive_icon_palette(colors),
        slide_width=prs.slide_width,
        slide_height=prs.slide_height,
        meta=Meta(
            spec_version=SPEC_VERSION,
            client=client,
            project=project,
            source_master=str(master_pptx_path),
            source_name=master_pptx_path.stem,
            extracted_at=datetime.now(timezone.utc).isoformat(timespec="seconds"),
            layouts_found=len(layouts),
            content_slides_ignored=len(prs.slides._sldIdLst),
        ),
    )
    spec.grid = derive_grid(layouts)
    return spec
