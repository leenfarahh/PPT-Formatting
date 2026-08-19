"""
Stage 3a - classifying a content slide into a layout archetype.

The submitted deck is rough: shapes are hand-placed, placeholders are often
unused, and slides carry no reliable indication of what kind of slide they
are. Before anything can be reformatted, each slide has to be read for what
it structurally *is* - a cover, a divider, a two-column comparison, a chart
slide - so it can be routed to the matching layout.

Classification works off an inventory of the slide's shapes plus a few
positional signals, and every decision carries the evidence that produced
it. That matters in review: a designer who disagrees with a routing can see
why the tool chose it, and `SlideClassification.archetype` can simply be
overridden without re-running anything.

Deliberately rule-based rather than learned: the rules are inspectable and
adjustable by the people who own the deck standards, and there is no
training corpus of labelled Prezlab slides to learn from.
"""
from __future__ import annotations

import re
from dataclasses import dataclass, field

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from . import archetypes
from . import rtl

# Text at or under this length reads as a heading rather than body copy.
HEADING_MAX_CHARS = 90
# A statement slide is a single short line set large.
STATEMENT_MAX_CHARS = 160
# Fraction of the slide a picture must cover to count as full-bleed.
FULL_BLEED_COVERAGE = 0.8
# Horizontal distance under which two blocks count as the same column.
COLUMN_TOLERANCE = 0.08

TITLE_PLACEHOLDERS = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}

_CLOSING_PATTERNS = re.compile(
    r"\b(thank\s*you|thanks|questions\?|q\s*&\s*a|get\s*in\s*touch|contact\s*us)\b"
    r"|شكرا|شكراً|أسئلة",
    re.IGNORECASE,
)
_AGENDA_PATTERNS = re.compile(
    r"\b(agenda|contents|table\s*of\s*contents|overview|what\s*we.?ll\s*cover)\b"
    r"|جدول\s*الأعمال|المحتويات",
    re.IGNORECASE,
)
# Straight and typographic quotes, plus Arabic quotation marks.
_QUOTE_CHARS = '"“”«»‘’„❝❞'
_ATTRIBUTION = re.compile(r"^\s*[—–-]{1,2}\s*\S")


@dataclass
class TextBlock:
    """One text-bearing shape, reduced to what classification needs."""
    text: str
    left: float
    top: float
    width: float
    height: float
    font_size_pt: float | None
    is_title_placeholder: bool
    shape: object = None

    @property
    def char_count(self) -> int:
        return len(self.text.strip())

    @property
    def is_heading_like(self) -> bool:
        return 0 < self.char_count <= HEADING_MAX_CHARS and "\n" not in self.text.strip()


@dataclass
class SlideFeatures:
    """Structural inventory of one content slide."""
    text_blocks: list = field(default_factory=list)
    pictures: list = field(default_factory=list)
    tables: list = field(default_factory=list)
    charts: list = field(default_factory=list)
    other_shapes: list = field(default_factory=list)
    max_picture_coverage: float = 0.0
    index: int = 0
    total: int = 1
    is_rtl: bool = False

    @property
    def has_content(self) -> bool:
        return bool(self.text_blocks or self.pictures or self.tables or self.charts)


@dataclass
class SlideClassification:
    archetype: str
    confidence: float
    evidence: list = field(default_factory=list)
    features: SlideFeatures | None = None


# --- feature extraction ---------------------------------------------------

def _largest_font_pt(shape) -> float | None:
    """The largest explicitly-sized run in a shape, if any size is set."""
    sizes = []
    if not shape.has_text_frame:
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                sizes.append(run.font.size.pt)
    return max(sizes) if sizes else None


def extract_features(slide, index: int, total: int, slide_w: int, slide_h: int) -> SlideFeatures:
    features = SlideFeatures(index=index, total=total)
    area = float(slide_w * slide_h) or 1.0

    def frac_geom(shape):
        if None in (shape.left, shape.top, shape.width, shape.height):
            return 0.0, 0.0, 0.0, 0.0
        return (shape.left / slide_w, shape.top / slide_h,
                shape.width / slide_w, shape.height / slide_h)

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            features.pictures.append(shape)
            if None not in (shape.width, shape.height):
                features.max_picture_coverage = max(
                    features.max_picture_coverage, (shape.width * shape.height) / area
                )
            continue
        if getattr(shape, "has_table", False) and shape.has_table:
            features.tables.append(shape)
            continue
        if getattr(shape, "has_chart", False) and shape.has_chart:
            features.charts.append(shape)
            continue
        if shape.has_text_frame and shape.text_frame.text.strip():
            left, top, width, height = frac_geom(shape)
            is_title = (
                shape.is_placeholder
                and shape.placeholder_format.type in TITLE_PLACEHOLDERS
            )
            features.text_blocks.append(TextBlock(
                text=shape.text_frame.text,
                left=left, top=top, width=width, height=height,
                font_size_pt=_largest_font_pt(shape),
                is_title_placeholder=is_title,
                shape=shape,
            ))
            continue
        features.other_shapes.append(shape)

    all_text = " ".join(b.text for b in features.text_blocks)
    features.is_rtl = rtl.is_rtl_text(all_text)
    return features


def identify_title(features: SlideFeatures) -> TextBlock | None:
    """
    Pick the block acting as the slide's title.

    A real title placeholder wins outright. Otherwise the heuristic favours
    a short block near the top, with an explicit large font as a tiebreak -
    which is how a hand-built rough slide signals a title without using a
    placeholder.
    """
    for block in features.text_blocks:
        if block.is_title_placeholder:
            return block

    candidates = [b for b in features.text_blocks if b.is_heading_like and b.top < 0.4]
    if not candidates:
        return None
    return max(candidates, key=lambda b: ((b.font_size_pt or 0), -b.top))


def _column_count(blocks: list) -> int:
    """How many distinct columns a set of blocks occupies."""
    if not blocks:
        return 1
    lefts = sorted(b.left for b in blocks)
    clusters = [lefts[0]]
    for value in lefts[1:]:
        if value - clusters[-1] > COLUMN_TOLERANCE:
            clusters.append(value)
    return len(clusters)


def _looks_like_quote(block: TextBlock) -> bool:
    text = block.text.strip()
    if not text:
        return False
    if text[0] in _QUOTE_CHARS and text.rstrip()[-1] in _QUOTE_CHARS + ".!?":
        return True
    lines = [l for l in text.split("\n") if l.strip()]
    # A quotation followed by an attribution line ("- Name, Title").
    return len(lines) >= 2 and bool(_ATTRIBUTION.match(lines[-1]))


# --- classification -------------------------------------------------------

def classify_slide(slide, index: int, total: int, slide_w: int, slide_h: int) -> SlideClassification:
    """
    Route one content slide to an archetype.

    Rules are ordered by how decisive their evidence is: an actual table on
    the slide beats any inference drawn from text lengths or positions.
    """
    f = extract_features(slide, index, total, slide_w, slide_h)
    ev: list = []

    if not f.has_content:
        return SlideClassification(archetypes.BLANK, 1.0, ["slide has no content shapes"], f)

    title = identify_title(f)
    body_blocks = [b for b in f.text_blocks if b is not title]
    all_text = " ".join(b.text for b in f.text_blocks)

    if title:
        ev.append(f"title identified: {title.text.strip()[:40]!r}")

    # -- decisive object types --------------------------------------------
    if f.tables:
        ev.append(f"{len(f.tables)} table(s) present")
        return SlideClassification(archetypes.TABLE, 0.95, ev, f)

    if f.charts:
        ev.append(f"{len(f.charts)} chart(s) present")
        return SlideClassification(archetypes.CHART, 0.95, ev, f)

    if f.max_picture_coverage >= FULL_BLEED_COVERAGE:
        ev.append(f"picture covers {f.max_picture_coverage:.0%} of the slide")
        return SlideClassification(archetypes.PICTURE_FULL, 0.9, ev, f)

    # -- position-anchored slides -----------------------------------------
    if _CLOSING_PATTERNS.search(all_text) and len(all_text) < 200:
        ev.append("closing language detected")
        return SlideClassification(archetypes.CLOSING, 0.85, ev, f)

    if index == 0 and len(f.text_blocks) <= 3 and not f.pictures:
        short = all(b.char_count <= STATEMENT_MAX_CHARS for b in f.text_blocks)
        if short:
            ev.append("first slide, few short text blocks")
            return SlideClassification(archetypes.TITLE_SLIDE, 0.8, ev, f)

    if _AGENDA_PATTERNS.search(all_text):
        ev.append("agenda/contents language detected")
        return SlideClassification(archetypes.TITLE_AND_CONTENT, 0.8, ev, f)

    # -- quotes and statements --------------------------------------------
    for block in f.text_blocks:
        if _looks_like_quote(block):
            ev.append("quotation marks or attribution line found")
            return SlideClassification(archetypes.QUOTE, 0.8, ev, f)

    if len(f.text_blocks) == 1 and not f.pictures:
        only = f.text_blocks[0]
        if only.char_count <= STATEMENT_MAX_CHARS:
            large = (only.font_size_pt or 0) >= 32
            ev.append(
                f"single short text block ({only.char_count} chars"
                + (", large type" if large else "") + ")"
            )
            # A single short line set large is a statement; the same line at
            # ordinary size in the upper half is a section divider.
            if large or only.top > 0.3:
                return SlideClassification(archetypes.BIG_STATEMENT, 0.75, ev, f)
            return SlideClassification(archetypes.TITLE_ONLY, 0.7, ev, f)

    if title and len(body_blocks) == 1 and body_blocks[0].char_count <= HEADING_MAX_CHARS \
            and not f.pictures:
        ev.append("title plus one short line")
        return SlideClassification(archetypes.SECTION_HEADER, 0.7, ev, f)

    # -- picture layouts ---------------------------------------------------
    if f.pictures and f.text_blocks:
        ev.append(f"{len(f.pictures)} picture(s) alongside text")
        return SlideClassification(archetypes.PICTURE_CAPTION, 0.75, ev, f)

    if f.pictures and not f.text_blocks:
        ev.append("pictures only, no text")
        return SlideClassification(archetypes.PICTURE_FULL, 0.7, ev, f)

    # -- multi-column content ---------------------------------------------
    columns = _column_count(body_blocks)
    ev.append(f"{len(body_blocks)} body block(s) across {columns} column(s)")

    if columns >= 2:
        headings = [b for b in body_blocks if b.is_heading_like]
        # Two columns of heading-plus-body pairs is a comparison, not just
        # a two-column split.
        if columns == 2 and len(body_blocks) >= 4 and len(headings) >= 2:
            ev.append("paired headings and bodies in two columns")
            return SlideClassification(archetypes.COMPARISON, 0.75, ev, f)
        if columns >= 3:
            return SlideClassification(archetypes.THREE_CONTENT, 0.75, ev, f)
        return SlideClassification(archetypes.TWO_CONTENT, 0.75, ev, f)

    if title and body_blocks:
        return SlideClassification(archetypes.TITLE_AND_CONTENT, 0.7, ev, f)

    if title and not body_blocks:
        return SlideClassification(archetypes.TITLE_ONLY, 0.7, ev, f)

    ev.append("no decisive signal; defaulting to title and content")
    return SlideClassification(archetypes.TITLE_AND_CONTENT, 0.4, ev, f)


def classify_deck(prs, slide_w: int, slide_h: int) -> list:
    """Classify every slide in a rough content deck."""
    slides = list(prs.slides)
    total = len(slides)
    return [
        classify_slide(slide, i, total, slide_w, slide_h)
        for i, slide in enumerate(slides)
    ]


# --- role assignment ------------------------------------------------------

# Archetypes whose layout is built around a dominant title. On these, text
# that would otherwise be treated as body copy belongs in the title, since
# that placeholder is what carries the large display type.
TITLE_DOMINANT = {
    archetypes.TITLE_SLIDE, archetypes.CLOSING, archetypes.SECTION_HEADER,
    archetypes.BIG_STATEMENT, archetypes.QUOTE, archetypes.TITLE_ONLY,
}


@dataclass
class TextSource:
    """
    Text destined for one placeholder.

    Usually a whole shape, but `paragraphs` lets a single shape be split
    across two placeholders - a quote slide where one text box holds both
    the quotation and its attribution needs exactly that.
    """
    block: TextBlock
    paragraphs: list | None = None

    @property
    def frame(self):
        return self.block.shape.text_frame

    def para_list(self) -> list:
        if self.paragraphs is not None:
            return self.paragraphs
        return list(self.frame.paragraphs)


def assign_roles(features: SlideFeatures, archetype: str) -> tuple:
    """
    Decide which text becomes the title and which becomes body copy, given
    what kind of slide this is.

    Classification alone isn't enough here. A rough "Thank you" slide has
    one text box and no title placeholder, so a purely positional rule
    treats it as body copy and drops it into a subtitle - technically
    consistent, visually wrong. Knowing the slide is a closing tells us
    that box *is* the title.

    Returns `(title_source, body_sources)`, either of which may be empty.
    """
    if features is None:
        return None, []

    title_block = identify_title(features)
    body = [b for b in features.text_blocks if b is not title_block]
    body.sort(key=lambda b: (round(b.top, 2), b.left))

    # On a title-dominant layout, promote the most prominent remaining
    # block rather than leaving the title placeholder empty.
    if title_block is None and archetype in TITLE_DOMINANT and body:
        promoted = max(body, key=lambda b: ((b.font_size_pt or 0), -b.top))
        body.remove(promoted)
        title_block = promoted

    title_source = TextSource(title_block) if title_block is not None else None
    body_sources = [TextSource(b) for b in body]

    # A quote and its attribution usually share one text box. Split them so
    # the quotation gets the display placeholder and the credit gets the
    # smaller one, instead of both being set at the same size.
    if archetype == archetypes.QUOTE and title_source is not None:
        paragraphs = [p for p in title_source.frame.paragraphs]
        non_empty = [p for p in paragraphs if p.text.strip()]
        if len(non_empty) >= 2 and _ATTRIBUTION.match(non_empty[-1].text.strip()):
            attribution = non_empty[-1]
            title_source.paragraphs = [p for p in paragraphs if p is not attribution]
            body_sources.insert(0, TextSource(title_source.block, [attribution]))

    return title_source, body_sources
