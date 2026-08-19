"""
Stage 3c - detailed formatting.

Runs over a rebuilt slide and applies typography, color, bilingual text
direction, grid alignment, and table/chart styling.

The governing rule here is that **placeholders are authoritative**. Content
mapped into a placeholder already sits at the layout's geometry and already
inherits the master's typography, so this stage must not pin sizes or
positions over it - doing so would recreate exactly the hardcoded overrides
Stage 2 works to remove. What placeholders do get is direction handling and
theme *references* (`+mn-lt`, `+mn-cs`), which set the right face per script
without breaking inheritance.

Shapes that could not be mapped into a placeholder are a different case:
they carry the rough deck's arbitrary geometry and literal colors, so they
get the full treatment - theme fonts, nearest-accent color remapping, and
snapping to the grid the master implies.
"""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from . import geometry, rtl
from .style_spec import StyleSpec, ACCENT_ROLES

# Sizes applied only to text that could not be placed in a placeholder;
# placeholder text inherits its size from the layout instead.
FALLBACK_TITLE_PT = Pt(28)
FALLBACK_BODY_PT = Pt(18)

TITLE_PLACEHOLDERS = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}

# Theme font reference tokens. Setting these rather than a literal typeface
# keeps a run following the theme.
MAJOR_LATIN, MINOR_LATIN = "+mj-lt", "+mn-lt"
MAJOR_CS, MINOR_CS = "+mj-cs", "+mn-cs"


# --- color ----------------------------------------------------------------

def _hex_to_rgb(hex_str: str) -> tuple:
    hex_str = hex_str.lstrip("#")
    return int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)


def _rgb_distance(a: tuple, b: tuple) -> float:
    """
    Perceptually weighted RGB distance.

    Plain Euclidean distance in RGB treats a green shift the same as a blue
    one, which the eye does not. These coefficients approximate relative
    sensitivity and are enough to stop, say, a mid-grey being pulled toward
    a saturated accent that happens to be numerically close.
    """
    return (
        2.0 * (a[0] - b[0]) ** 2
        + 4.0 * (a[1] - b[1]) ** 2
        + 3.0 * (a[2] - b[2]) ** 2
    ) ** 0.5


def nearest_accent_role(rgb_hex: str, spec: StyleSpec) -> str | None:
    """Map an arbitrary RGB to the closest theme accent role."""
    target = _hex_to_rgb(rgb_hex)
    best_role, best_dist = None, float("inf")
    for role in ACCENT_ROLES:
        if role not in spec.theme.colors:
            continue
        dist = _rgb_distance(target, _hex_to_rgb(spec.theme.colors[role]))
        if dist < best_dist:
            best_role, best_dist = role, dist
    return best_role


def apply_color_mapping(slide, spec: StyleSpec) -> list:
    """
    Remap literal fill colors on unmapped shapes to the nearest theme accent.

    Placeholders are skipped: their fill comes from the layout, and
    overwriting it would pin a color the theme should control.
    """
    notes: list = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            continue
        if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART):
            continue
        try:
            fill = shape.fill
        except (AttributeError, TypeError):
            continue
        if fill is None or fill.type is None or int(fill.type) != 1:   # 1 == solid
            continue
        try:
            current = str(fill.fore_color.rgb)
        except (AttributeError, TypeError):
            continue    # already a theme reference, or not a literal RGB

        role = nearest_accent_role(current, spec)
        if not role:
            continue
        new_hex = spec.theme.colors[role]
        if new_hex.upper() != current.upper():
            fill.fore_color.rgb = RGBColor.from_string(new_hex)
            notes.append(f"{shape.name}: #{current} -> #{new_hex} ({role})")
    return notes


# --- typography and direction --------------------------------------------

def _is_title(shape) -> bool:
    return shape.is_placeholder and shape.placeholder_format.type in TITLE_PLACEHOLDERS


def _apply_theme_fonts(paragraph, is_title: bool) -> None:
    """
    Point every run at the theme's fonts rather than a literal typeface.

    Both the Latin and complex-script faces are set on every run, because a
    bilingual line needs each: PowerPoint picks per character, so an Arabic
    word inside an English sentence resolves through `+mn-cs` while the
    English resolves through `+mn-lt`.
    """
    latin = MAJOR_LATIN if is_title else MINOR_LATIN
    cs = MAJOR_CS if is_title else MINOR_CS
    for run in paragraph.runs:
        run.font.name = latin
        rtl.set_complex_script_font(run, cs)


def apply_typography(slide, spec: StyleSpec) -> dict:
    """
    Normalize text frames, set direction per paragraph, and point runs at
    the theme fonts.

    Sizes are pinned only on non-placeholder text. Word wrap is forced on
    everywhere: a plain text box defaults to `wrap="none"` with autofit, so
    a single long line silently grows past the shape - and off the slide -
    instead of respecting the box it was drawn in.
    """
    stats = {"rtl_paragraphs": 0, "frames": 0}

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        stats["frames"] += 1
        is_title = _is_title(shape)
        frame = shape.text_frame
        frame.word_wrap = True
        if frame.auto_size != MSO_AUTO_SIZE.NONE:
            frame.auto_size = MSO_AUTO_SIZE.NONE

        for para in frame.paragraphs:
            text = "".join(run.text for run in para.runs)
            if text.strip():
                is_rtl = rtl.is_rtl_text(text)
                rtl.set_paragraph_direction(para, is_rtl)
                if is_rtl:
                    rtl.align_for_direction(para, True)
                    stats["rtl_paragraphs"] += 1

            _apply_theme_fonts(para, is_title)

            if not shape.is_placeholder:
                # Unmapped text has no layout to inherit from, so it needs
                # concrete values.
                if para.alignment is None:
                    para.alignment = PP_ALIGN.LEFT
                para.line_spacing = 1.15
                for run in para.runs:
                    if run.font.size is None:
                        run.font.size = FALLBACK_TITLE_PT if is_title else FALLBACK_BODY_PT

    return stats


# --- grid -----------------------------------------------------------------

def _nearest(value: float, guides: list) -> float:
    return min(guides, key=lambda g: abs(g - value)) if guides else value


def apply_grid_alignment(slide, spec: StyleSpec) -> list:
    """
    Snap unmapped shapes to the grid inferred from the master.

    Placeholders are left alone - they are already exactly where the layout
    puts them, and snapping them would move content off the design. Only
    shapes carried over from the rough deck are adjusted, and only their
    position: resizing them could clip their content.

    This is a snap, not a constraint solver: it will not resolve an overlap
    that the rough deck already had. What it does guarantee is that it never
    *creates* one. A snap that would push a shape onto a placeholder or onto
    a shape already snapped is abandoned and the shape left where it was,
    because a few millimetres off the guide is a far smaller defect than
    content buried under content.
    """
    grid = spec.grid
    moved: list = []
    sw, sh = spec.slide_width, spec.slide_height
    if not grid.column_guides or not grid.row_guides:
        grid.compute_guides()

    max_left = 1.0 - grid.margin_right_frac
    max_top = 1.0 - grid.margin_bottom_frac

    # Placeholders are fixed points: they are where the layout put them and
    # nothing may be snapped on top of them.
    occupied = [
        box for box in (geometry.rect(s) for s in slide.shapes if s.is_placeholder) if box
    ]
    movable = [
        s for s in slide.shapes
        if not s.is_placeholder and s.left is not None and s.top is not None
    ]
    # A shape that stays put still occupies its box, so seed the set with
    # every candidate and swap the entry when one actually moves.
    boxes = {id(s): geometry.rect(s) for s in movable}
    occupied.extend(box for box in boxes.values() if box)

    for shape in movable:
        left_frac = shape.left / sw
        top_frac = shape.top / sh
        width_frac = (shape.width or 0) / sw
        height_frac = (shape.height or 0) / sh

        new_left = _nearest(left_frac, grid.column_guides)
        new_top = _nearest(top_frac, grid.row_guides)

        # Keep the shape inside the margins after snapping.
        new_left = max(grid.margin_left_frac, min(new_left, max(0.0, max_left - width_frac)))
        new_top = max(grid.margin_top_frac, min(new_top, max(0.0, max_top - height_frac)))

        if abs(new_left - left_frac) <= 1e-6 and abs(new_top - top_frac) <= 1e-6:
            continue

        current = boxes[id(shape)]
        target = None
        if current is not None:
            left_emu, top_emu = int(new_left * sw), int(new_top * sh)
            target = (
                left_emu, top_emu,
                left_emu + (shape.width or 0), top_emu + (shape.height or 0),
            )
            others = [b for b in occupied if b is not current]
            if geometry.collides(target, others):
                continue

        shape.left = Emu(int(new_left * sw))
        shape.top = Emu(int(new_top * sh))
        moved.append(shape.name)

        if current is not None:
            occupied.remove(current)
            occupied.append(target)
            boxes[id(shape)] = target

    return moved


# --- tables and charts ----------------------------------------------------

def format_table(table_shape, spec: StyleSpec) -> None:
    """Apply the brand's table styling: header shading, borders, fonts."""
    style = spec.table_style
    table = table_shape.table

    for r_idx, row in enumerate(table.rows):
        is_header = r_idx == 0
        for cell in row.cells:
            if is_header and style.header_fill:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(style.header_fill)
            elif style.banded_fill and r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(style.banded_fill)

            cell.margin_left = Emu(style.cell_padding_emu)
            cell.margin_right = Emu(style.cell_padding_emu)

            for para in cell.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs)
                if text.strip():
                    is_rtl = rtl.is_rtl_text(text)
                    rtl.set_paragraph_direction(para, is_rtl)
                    if is_rtl:
                        rtl.align_for_direction(para, True)
                for run in para.runs:
                    run.font.name = MINOR_LATIN
                    rtl.set_complex_script_font(run, MINOR_CS)
                    run.font.size = Pt(
                        style.header_font_size_pt if is_header else style.body_font_size_pt
                    )
                    if is_header:
                        run.font.bold = style.header_bold
                        if style.header_font_color:
                            run.font.color.rgb = RGBColor.from_string(style.header_font_color)
                    elif style.body_font_color:
                        run.font.color.rgb = RGBColor.from_string(style.body_font_color)


def format_chart(chart_shape, spec: StyleSpec) -> list:
    """
    Recolor a native chart's series to the brand palette and restyle its
    text.

    Chart XML varies a lot by chart type, and a series that carries no
    fill element behaves differently from one that does, so each step is
    guarded: a chart that resists styling is reported rather than allowed
    to fail the deck.
    """
    notes: list = []
    style = spec.chart_style
    if not style.series_colors:
        return notes

    try:
        chart = chart_shape.chart
    except (AttributeError, ValueError):
        return notes

    try:
        chart.font.size = Pt(style.font_size_pt)
        chart.font.name = MINOR_LATIN
    except (AttributeError, ValueError):
        pass

    try:
        for i, series in enumerate(chart.series):
            color = style.series_colors[i % len(style.series_colors)]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor.from_string(color)
        notes.append(f"recolored {len(list(chart.series))} chart series")
    except (AttributeError, ValueError, NotImplementedError) as exc:
        notes.append(f"could not recolor series on '{chart_shape.name}': {exc}")

    try:
        if chart.has_legend:
            chart.legend.include_in_layout = False
    except (AttributeError, ValueError):
        pass

    return notes


# --- entry point ----------------------------------------------------------

def format_slide(slide, spec: StyleSpec) -> dict:
    """Run the full Stage 3 formatting pass over one rebuilt slide."""
    typography = apply_typography(slide, spec)
    color_notes = apply_color_mapping(slide, spec)
    moved = apply_grid_alignment(slide, spec)

    tables = charts = 0
    chart_notes: list = []
    for shape in slide.shapes:
        if getattr(shape, "has_table", False) and shape.has_table:
            format_table(shape, spec)
            tables += 1
        elif getattr(shape, "has_chart", False) and shape.has_chart:
            chart_notes.extend(format_chart(shape, spec))
            charts += 1

    return {
        "color_remaps": color_notes,
        "shapes_snapped": moved,
        "tables_formatted": tables,
        "charts_formatted": charts,
        "chart_notes": chart_notes,
        "rtl_paragraphs": typography["rtl_paragraphs"],
        "text_frames": typography["frames"],
    }
