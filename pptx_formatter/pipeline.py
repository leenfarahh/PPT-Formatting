"""
End-to-end orchestration.

Three stages, each independently runnable so its output can be inspected
and diffed on its own:

    Stage 1  ingest_master()    master.pptx        -> Style Spec (+ banked)
    Stage 2  build_template()   Style Spec         -> restyled deck of layouts
    Stage 3  apply_to_deck()    rough content deck -> formatted deck

`run_pipeline()` chains all three. `format_with_banked_spec()` skips Stage 1
for a repeat client by starting from the archived spec, which is the payoff
of banking specs in the first place.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from . import archetypes, classifier, qa
from .bank import TemplateBank
from .extraction import extract_style_spec
from .formatting import format_slide
from .layout_generator import generate_master_layouts, find_layout
from .slide_copy import rebuild_slide
from .style_spec import StyleSpec


# --- Stage 1 --------------------------------------------------------------

def ingest_master(
    master_pptx_path: str | Path,
    bank: TemplateBank | None = None,
    client: str | None = None,
    project: str | None = None,
) -> tuple[StyleSpec, str | None]:
    """
    Read a submitted master into a Style Spec and archive it.

    The submission may be an empty master or a fully populated deck; either
    way only its masters and layouts are read and its content slides are
    ignored. Returns the spec and its bank entry id (None when unbanked).
    """
    master_pptx_path = Path(master_pptx_path)

    if bank is None:
        return extract_style_spec(master_pptx_path, client=client, project=project), None

    # Allocate the entry first so extracted assets land straight in the bank
    # rather than in a temp directory that then has to be copied in.
    entry_id = bank.allocate(client, master_pptx_path.stem)
    spec = extract_style_spec(
        master_pptx_path,
        asset_dir=bank.asset_dir(entry_id),
        client=client,
        project=project,
    )
    bank.save(spec, entry_id, master_pptx=master_pptx_path)
    return spec, entry_id


# --- Stage 2 --------------------------------------------------------------

def build_template(
    spec: StyleSpec,
    base_pptx: str | Path,
    bank: TemplateBank | None = None,
    exclude_entry: str | None = None,
    wanted_archetypes: list | None = None,
):
    """
    Turn a Style Spec into a deck carrying a full set of restyled layouts.

    `base_pptx` is the designer's own master, which is what keeps their
    authored layouts intact rather than approximating them.
    """
    return generate_master_layouts(
        spec, base_pptx, bank=bank,
        wanted_archetypes=wanted_archetypes, exclude_entry=exclude_entry,
    )


# --- Stage 3 --------------------------------------------------------------

def apply_to_deck(
    out_prs,
    rough_content_pptx_path: str | Path,
    spec: StyleSpec,
    layout_archetypes: dict,
    layout_overrides: dict | None = None,
) -> dict:
    """
    Reformat every slide of a rough deck onto the matching layout.

    Each slide is classified structurally, routed to the layout carrying
    that archetype (falling back along the archetype's fallback chain when
    the deck has no exact match), rebuilt onto it, then formatted.

    `layout_overrides` maps a slide index to an archetype, letting a
    designer correct a routing without touching the classifier.
    """
    rough = Presentation(str(rough_content_pptx_path))
    overrides = layout_overrides or {}

    # archetype -> layout name, for routing.
    by_archetype: dict = {}
    for name, archetype in layout_archetypes.items():
        by_archetype.setdefault(archetype, name)

    background_hex = spec.theme.colors.get("lt1", "FFFFFF")
    slides_report = []
    all_warnings: list = []
    all_qa: list = []

    classifications = classifier.classify_deck(rough, spec.slide_width, spec.slide_height)

    for index, (source_slide, classification) in enumerate(zip(rough.slides, classifications)):
        archetype = overrides.get(index, classification.archetype)
        if index in overrides:
            classification.evidence.append(f"overridden to '{archetype}' by caller")

        layout_name, resolved = _resolve_layout(archetype, by_archetype)
        layout = find_layout(out_prs, layout_name) if layout_name else None
        if layout is None:
            # Should not happen once Stage 2 has run, but never fail a deck
            # over a routing miss - put the slide somewhere sensible.
            layout = out_prs.slide_masters[0].slide_layouts[0]
            all_warnings.append(
                f"slide {index + 1}: no layout for '{archetype}'; used '{layout.name}'"
            )

        new_slide = out_prs.slides.add_slide(layout)
        copy_report = rebuild_slide(new_slide, source_slide, classification, spec, layout)
        format_report = format_slide(new_slide, spec)
        issues = qa.check_slide(new_slide, spec, background_hex)

        all_warnings.extend(f"slide {index + 1}: {w}" for w in copy_report["warnings"])
        all_qa.extend(f"slide {index + 1}: {i}" for i in issues)

        slides_report.append({
            "slide": index + 1,
            "archetype": archetype,
            "resolved_archetype": resolved,
            "confidence": classification.confidence,
            "evidence": classification.evidence,
            "layout": layout.name,
            "mapped": copy_report["mapped"],
            "carried_over": copy_report["carried_over"],
            "dropped": copy_report["dropped"],
            "unused_placeholders": copy_report["unused_placeholders"],
            "warnings": copy_report["warnings"],
            "formatting": format_report,
            "qa_issues": issues,
        })

    return {
        "slides": slides_report,
        "slides_processed": len(slides_report),
        "warnings": all_warnings,
        "qa_issues": all_qa,
    }


def _resolve_layout(archetype: str, by_archetype: dict) -> tuple:
    """
    Find a layout for an archetype, walking the fallback chain if needed.

    Returns the layout name and the archetype actually used, so a report
    can show where a slide was routed when its ideal layout didn't exist.
    """
    if archetype in by_archetype:
        return by_archetype[archetype], archetype
    for candidate in archetypes.FALLBACK_CHAIN.get(archetype, []):
        if candidate in by_archetype:
            return by_archetype[candidate], candidate
    if archetypes.TITLE_AND_CONTENT in by_archetype:
        return by_archetype[archetypes.TITLE_AND_CONTENT], archetypes.TITLE_AND_CONTENT
    return (next(iter(by_archetype.values())), None) if by_archetype else (None, None)


# --- the whole thing ------------------------------------------------------

def run_pipeline(
    master_pptx_path: str | Path,
    rough_content_pptx_path: str | Path,
    output_pptx_path: str | Path,
    bank: TemplateBank | None = None,
    client: str | None = None,
    project: str | None = None,
    layout_overrides: dict | None = None,
) -> dict:
    """
    Stages 1-3: a submitted master plus a rough deck becomes a formatted,
    fully editable `.pptx`.

    Returns a report covering the extracted spec, where each layout came
    from, how each slide was classified and routed, and any QA findings.
    """
    output_pptx_path = Path(output_pptx_path)

    spec, entry_id = ingest_master(master_pptx_path, bank=bank, client=client, project=project)

    out_prs, stage2 = build_template(
        spec, master_pptx_path, bank=bank, exclude_entry=entry_id
    )

    stage3 = apply_to_deck(
        out_prs, rough_content_pptx_path, spec,
        stage2["layout_archetypes"], layout_overrides=layout_overrides,
    )

    output_pptx_path.parent.mkdir(parents=True, exist_ok=True)
    out_prs.save(str(output_pptx_path))

    return {
        "output_path": str(output_pptx_path),
        "bank_entry": entry_id,
        "style_spec": spec.to_dict(),
        "master": {
            "content_slides_ignored": spec.meta.content_slides_ignored,
            "layouts_found": spec.meta.layouts_found,
        },
        "layouts": {
            "designer": stage2["designer_layouts"],
            "from_bank": stage2["bank_layouts"],
            "generated": stage2["generated_layouts"],
            "available": stage2["available_layouts"],
            "inheritance": stage2["inheritance"],
        },
        "slides": stage3["slides"],
        "slides_processed": stage3["slides_processed"],
        "warnings": stage3["warnings"],
        "qa_issues": stage3["qa_issues"],
    }


def format_with_banked_spec(
    client: str,
    rough_content_pptx_path: str | Path,
    output_pptx_path: str | Path,
    bank: TemplateBank,
    layout_overrides: dict | None = None,
) -> dict:
    """
    Format a deck for a returning client without a master submission.

    Stage 1 is skipped entirely: the archived spec and master for the
    client's most recent entry are used instead. Raises LookupError when
    the bank holds nothing for that client.
    """
    entry = bank.latest_entry_for_client(client)
    if entry is None:
        raise LookupError(f"No banked Style Spec for client '{client}'")
    if not bank.has_master(entry.entry_id):
        raise LookupError(
            f"Bank entry '{entry.entry_id}' has no archived master to build on; "
            "re-submit the master once to populate it"
        )

    spec = bank.load(entry.entry_id)
    out_prs, stage2 = build_template(
        spec, bank.master_path(entry.entry_id), bank=bank, exclude_entry=None
    )
    stage3 = apply_to_deck(
        out_prs, rough_content_pptx_path, spec,
        stage2["layout_archetypes"], layout_overrides=layout_overrides,
    )

    output_pptx_path = Path(output_pptx_path)
    output_pptx_path.parent.mkdir(parents=True, exist_ok=True)
    out_prs.save(str(output_pptx_path))

    return {
        "output_path": str(output_pptx_path),
        "bank_entry": entry.entry_id,
        "stage_1_skipped": True,
        "layouts": {
            "designer": stage2["designer_layouts"],
            "from_bank": stage2["bank_layouts"],
            "generated": stage2["generated_layouts"],
            "available": stage2["available_layouts"],
            "inheritance": stage2["inheritance"],
        },
        "slides": stage3["slides"],
        "slides_processed": stage3["slides_processed"],
        "warnings": stage3["warnings"],
        "qa_issues": stage3["qa_issues"],
    }
