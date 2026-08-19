"""
Stage 2 - Master layout generation.

The output deck is built on the designer's *own* master, with its content
slides stripped. That single decision is what preserves their layouts
natively: their placeholder geometry, text styles, color map and theme part
all survive because we never left their file. Earlier the tool restyled a
generic bundled template and threw the designer's layouts away.

From that base the stage does four things:

1.  **Fills gaps.** Any canonical archetype the designer didn't author is
    sourced from the Template Bank - the structurally closest banked layout
    wins - or, when the bank has nothing, synthesized from the submission's
    own grid. Both arrive as real `sldLayout` parts via `layout_builder`.

2.  **Restyles every layout**, designer-authored and template-sourced
    alike: theme colors and fonts substituted into the theme part, so
    anything referencing the theme follows automatically.

3.  **Propagates logo and footer.** The logo is placed once on the master
    and inherited; layouts its placement rule excludes suppress master
    shapes instead of carrying a duplicate copy.

4.  **Verifies placeholder inheritance.** Hardcoded fonts and colors that
    merely restate a theme value are rewritten as theme references
    (`+mn-lt`, `<a:schemeClr val="tx1"/>`), so layouts follow the master
    rather than pinning values that a later theme change won't reach.
"""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from . import archetypes
from .style_spec import StyleSpec, LayoutSpec, THEME_COLOR_ROLES, SOURCE_DESIGNER
from .extraction import get_theme_part, A_NS, P_NS
from . import layout_builder
from .bank import TemplateBank

# In a layout or slide, theme colors are addressed through the master's
# color map, not by their theme-part role names. These are the mapped
# equivalents used inside <a:schemeClr>.
CLR_MAP = {"dk1": "tx1", "lt1": "bg1", "dk2": "tx2", "lt2": "bg2"}

# Theme font reference tokens understood by PowerPoint.
FONT_TOKENS = {
    ("major", "latin"): "+mj-lt", ("minor", "latin"): "+mn-lt",
    ("major", "ea"): "+mj-ea", ("minor", "ea"): "+mn-ea",
    ("major", "cs"): "+mj-cs", ("minor", "cs"): "+mn-cs",
}


# --- theme rewriting ------------------------------------------------------

def _write_theme_colors(theme_root: etree._Element, colors: dict) -> None:
    scheme = theme_root.find(".//a:clrScheme", A_NS)
    if scheme is None:
        return
    for role in THEME_COLOR_ROLES:
        if role not in colors:
            continue
        role_el = scheme.find(f"a:{role}", A_NS)
        if role_el is None:
            continue
        srgb = role_el.find("a:srgbClr", A_NS)
        if srgb is not None:
            srgb.set("val", colors[role].upper().lstrip("#"))
            continue
        # dk1/lt1 commonly use sysClr; replace with an explicit srgbClr so
        # the brand color survives rather than following the OS theme.
        sys_clr = role_el.find("a:sysClr", A_NS)
        if sys_clr is not None:
            role_el.remove(sys_clr)
        new_srgb = etree.SubElement(role_el, f"{{{A_NS['a']}}}srgbClr")
        new_srgb.set("val", colors[role].upper().lstrip("#"))


def _write_theme_fonts(theme_root: etree._Element, fonts) -> None:
    """
    Write the per-script typefaces into the theme.

    All three scripts are written, including the complex-script face even
    when the theme declared none of its own. Formatting points Arabic runs
    at `+mn-cs`, and that token resolves through this entry: leaving it
    empty would send Arabic text to whatever font the viewer's machine
    happens to substitute. `cs_for()` falls back to the Latin face, so the
    worst case is the same font the deck already uses.
    """
    scheme = theme_root.find(".//a:fontScheme", A_NS)
    if scheme is None:
        return
    for prefix, is_major in (("major", True), ("minor", False)):
        for script, value in (
            ("latin", fonts.latin_for(is_major)),
            ("ea", fonts.ea_for(is_major)),
            ("cs", fonts.cs_for(is_major)),
        ):
            if not value:
                continue
            el = scheme.find(f"a:{prefix}Font/a:{script}", A_NS)
            if el is not None:
                el.set("typeface", value)


def apply_theme(prs: Presentation, spec: StyleSpec) -> None:
    """Substitute the brand's colors and per-script fonts into the theme part."""
    theme_part = get_theme_part(prs)
    theme_root = etree.fromstring(theme_part.blob)
    _write_theme_colors(theme_root, spec.theme.colors)
    _write_theme_fonts(theme_root, spec.theme.fonts)
    theme_part._blob = etree.tostring(
        theme_root, xml_declaration=True, encoding="UTF-8", standalone=True
    )


# --- inheritance repair ---------------------------------------------------

def normalize_theme_references(element: etree._Element, spec: StyleSpec) -> list:
    """
    Rewrite hardcoded values that merely restate the theme as references
    to it, and report what changed.

    A layout that pins `typeface="Calibri"` and `srgbClr val="4F81BD"` looks
    identical today but stops following the master the moment the theme
    changes - which is exactly what happens when this tool restyles a deck
    for the next client. Converting those to `+mn-lt` and
    `<a:schemeClr val="accent1"/>` is what "inherits from the master"
    actually means at the XML level.

    Values that don't correspond to any theme entry are left alone: those
    are deliberate design decisions, not accidental overrides.
    """
    changes: list = []
    fonts = spec.theme.fonts
    colors = spec.theme.colors

    font_lookup = {}
    for is_major, prefix in ((True, "major"), (False, "minor")):
        for script, value in (
            ("latin", fonts.latin_for(is_major)),
            ("ea", fonts.ea_for(is_major)),
            ("cs", fonts.cs_for(is_major)),
        ):
            if value:
                font_lookup.setdefault((script, value.lower()), FONT_TOKENS[(prefix, script)])

    for script in ("latin", "ea", "cs"):
        for el in element.iter(f"{{{A_NS['a']}}}{script}"):
            typeface = el.get("typeface")
            if not typeface or typeface.startswith("+"):
                continue
            token = font_lookup.get((script, typeface.lower()))
            if token:
                el.set("typeface", token)
                changes.append(f"font '{typeface}' -> theme reference {token}")

    hex_to_role = {v.upper(): k for k, v in colors.items()}
    for srgb in list(element.iter(f"{{{A_NS['a']}}}srgbClr")):
        val = (srgb.get("val") or "").upper()
        role = hex_to_role.get(val)
        if not role:
            continue
        parent = srgb.getparent()
        if parent is None:
            continue
        scheme_clr = parent.makeelement(f"{{{A_NS['a']}}}schemeClr", {})
        scheme_clr.set("val", CLR_MAP.get(role, role))
        # Carry over child transforms (alpha, tint, shade) so the visual
        # result is unchanged; only the color's source becomes the theme.
        for child in srgb:
            scheme_clr.append(child)
        parent.replace(srgb, scheme_clr)
        changes.append(f"color #{val} -> theme role {CLR_MAP.get(role, role)}")

    return changes


def verify_placeholder_inheritance(prs: Presentation, spec: StyleSpec) -> dict:
    """
    Walk every layout, rewriting theme-equivalent hardcoded values into
    theme references and confirming each placeholder resolves against the
    master.

    Returns a report of what was repaired and what could not be.
    """
    repaired: list = []
    orphans: list = []
    master = prs.slide_masters[0]
    master_ph_types = set()
    for shape in master.placeholders:
        ph = shape._element.find(".//p:nvSpPr/p:nvPr/p:ph", P_NS)
        if ph is not None:
            master_ph_types.add(ph.get("type") or "body")

    for layout in master.slide_layouts:
        changes = normalize_theme_references(layout._element, spec)
        if changes:
            repaired.append({"layout": layout.name, "changes": changes})

        for shape in layout.placeholders:
            ph = shape._element.find(".//p:nvSpPr/p:nvPr/p:ph", P_NS)
            ph_type = (ph.get("type") if ph is not None else None) or "body"
            # body/obj/pic/tbl/chart all inherit from the master's body
            # placeholder; title and ctrTitle from its title.
            base = {
                "ctrTitle": "title", "subTitle": "body", "obj": "body",
                "pic": "body", "tbl": "body", "chart": "body", "dgm": "body",
                "media": "body", "clipArt": "body", "sldImg": "body",
            }.get(ph_type, ph_type)
            if base not in master_ph_types and base not in ("dt", "ftr", "sldNum"):
                orphans.append(
                    f"'{layout.name}' placeholder '{shape.name}' ({ph_type}) has no "
                    f"matching '{base}' placeholder on the master to inherit from"
                )

    # The master itself can pin theme values too, and everything inherits
    # from it, so it is worth normalizing as well.
    master_changes = normalize_theme_references(master._element, spec)
    if master_changes:
        repaired.append({"layout": "<slide master>", "changes": master_changes})

    return {"repaired": repaired, "orphans": orphans}


# --- logo, footer, background --------------------------------------------

def apply_logo(prs: Presentation, spec: StyleSpec) -> str | None:
    """
    Place the brand logo once on the slide master, so every layout that
    shows master shapes inherits it.

    Existing master pictures matching the logo's footprint are removed
    first, which keeps re-running the generator idempotent instead of
    stacking copies.
    """
    logo = spec.brand.logo
    if not logo.present or not Path(logo.asset_path).exists():
        return None

    master = prs.slide_masters[0]
    sw, sh = prs.slide_width, prs.slide_height
    left = int(logo.left_frac * sw)
    top = int(logo.top_frac * sh)
    width = int(logo.width_frac * sw)
    height = int(logo.height_frac * sh)

    for shape in list(master.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        # Same footprint (within a 2% tolerance) means this is the logo we
        # are about to re-place, not unrelated master artwork.
        if (abs(shape.left - left) < sw * 0.02 and abs(shape.top - top) < sh * 0.02
                and abs(shape.width - width) < sw * 0.02):
            shape._element.getparent().remove(shape._element)

    layout_builder.add_picture_to(
        master.shapes, str(logo.asset_path), left, top, width, height, name="Brand Logo"
    )
    return "Brand Logo"


def apply_footer_text(prs: Presentation, spec: StyleSpec) -> None:
    """Write the brand's footer text into every footer placeholder."""
    footer = spec.brand.footer
    if not footer.text:
        return
    for master in prs.slide_masters:
        containers = [master] + list(master.slide_layouts)
        for container in containers:
            for ph in container.placeholders:
                if "FOOTER" in str(ph.placeholder_format.type or "") and ph.has_text_frame:
                    ph.text_frame.text = footer.text


def apply_master_background(prs: Presentation, spec: StyleSpec) -> None:
    """Apply the extracted master background, if it declared one."""
    bg = spec.master_background
    if bg.kind == "inherit":
        return
    master = prs.slide_masters[0]
    cSld = master._element.find("p:cSld", P_NS)
    if cSld is None:
        return
    for existing in cSld.findall("p:bg", P_NS):
        cSld.remove(existing)

    if bg.kind == "image" and bg.asset_path and Path(bg.asset_path).exists():
        layout_builder._apply_image_background(master.part, bg.asset_path)
        return
    xml = layout_builder._background_xml(bg)
    if xml:
        cSld.insert(0, etree.fromstring(xml.encode("utf-8")))


def _suppress_master_shapes(prs: Presentation, spec: StyleSpec, layout_archetypes: dict) -> list:
    """
    Set `showMasterSp="0"` on layouts the logo rule excludes.

    Suppressing inherited master shapes is how PowerPoint itself keeps a
    logo off a cover slide; it beats stamping a second copy of the logo
    onto each layout that wants one.
    """
    logo = spec.brand.logo
    if not logo.present:
        return []
    suppressed = []
    for layout in prs.slide_masters[0].slide_layouts:
        archetype = layout_archetypes.get(layout.name)
        if archetype is None:
            continue
        if logo.appears_on(archetype):
            if layout._element.get("showMasterSp") == "0":
                del layout._element.attrib["showMasterSp"]
        else:
            layout._element.set("showMasterSp", "0")
            suppressed.append(layout.name)
    return suppressed


# --- slide stripping ------------------------------------------------------

def strip_content_slides(prs: Presentation) -> int:
    """
    Remove every content slide from the master submission.

    The designer may submit an empty master or a fully populated deck; the
    layouts are what we want either way, so the slides go.
    """
    sldIdLst = prs.slides._sldIdLst
    removed = 0
    for sldId in list(sldIdLst):
        prs.part.drop_rel(sldId.rId)
        sldIdLst.remove(sldId)
        removed += 1
    return removed


# --- entry point ----------------------------------------------------------

def generate_master_layouts(
    spec: StyleSpec,
    base_pptx: str | Path,
    bank: TemplateBank | None = None,
    wanted_archetypes: list | None = None,
    exclude_entry: str | None = None,
) -> tuple[Presentation, dict]:
    """
    Stage 2 entry point: a Style Spec plus the master it came from becomes a
    restyled deck carrying a full set of layouts.

    `base_pptx` is the designer's master (or a banked one), used as the
    foundation so their authored layouts survive intact. Returns the
    presentation and a report describing where every layout came from.
    """
    prs = Presentation(str(base_pptx))
    report: dict = {
        "slides_stripped": strip_content_slides(prs),
        "designer_layouts": [],
        "bank_layouts": [],
        "generated_layouts": [],
        "layout_sources": {},
        "suppressed_master_shapes": [],
    }

    apply_theme(prs, spec)
    apply_master_background(prs, spec)
    apply_logo(prs, spec)

    # Which archetypes did the designer actually author?
    present: dict = {}
    for layout_spec in spec.layouts:
        if layout_spec.source == SOURCE_DESIGNER:
            present.setdefault(layout_spec.archetype, layout_spec.name)

    # layout name -> archetype, for the whole deck. Used to decide per-layout
    # logo suppression and to route content slides in Stage 3.
    layout_archetypes = {name: arch for arch, name in present.items()}
    for layout_spec in spec.layouts:
        layout_archetypes.setdefault(layout_spec.name, layout_spec.archetype)

    report["designer_layouts"] = [
        {"name": name, "archetype": arch} for arch, name in sorted(present.items())
    ]
    for arch, name in present.items():
        report["layout_sources"][name] = SOURCE_DESIGNER

    # Fill the gaps.
    targets = wanted_archetypes if wanted_archetypes is not None else archetypes.ALL_ARCHETYPES
    for archetype in targets:
        if archetype in present:
            continue

        chosen: LayoutSpec | None = None
        score = None
        if bank is not None:
            hit = bank.select_layout(archetype, spec, exclude_entry=exclude_entry)
            if hit:
                chosen, score = hit

        if chosen is None:
            chosen = layout_builder.generated_layout_spec(archetype, spec)

        # Name generated/banked layouts by archetype so the deck reads
        # consistently in PowerPoint's layout gallery, and so Stage 3 can
        # find them by name.
        chosen.name = _unique_layout_name(prs, archetypes.label_for(archetype))
        show_master = spec.brand.logo.appears_on(archetype) if spec.brand.logo.present else True
        layout_builder.add_layout(prs, chosen, show_master_shapes=show_master)

        layout_archetypes[chosen.name] = archetype
        report["layout_sources"][chosen.name] = chosen.source
        entry = {"name": chosen.name, "archetype": archetype, "source": chosen.source}
        if chosen.is_banked:
            entry["match_score"] = round(score, 2) if score is not None else None
            report["bank_layouts"].append(entry)
        else:
            report["generated_layouts"].append(entry)

    apply_footer_text(prs, spec)
    report["suppressed_master_shapes"] = _suppress_master_shapes(prs, spec, layout_archetypes)
    report["inheritance"] = verify_placeholder_inheritance(prs, spec)
    report["layout_archetypes"] = layout_archetypes
    report["available_layouts"] = list_available_layouts(prs)
    return prs, report


def _unique_layout_name(prs: Presentation, base: str) -> str:
    existing = {l.name for m in prs.slide_masters for l in m.slide_layouts}
    if base not in existing:
        return base
    n = 2
    while f"{base} {n}" in existing:
        n += 1
    return f"{base} {n}"


def list_available_layouts(prs: Presentation) -> list:
    return [layout.name for master in prs.slide_masters for layout in master.slide_layouts]


def find_layout(prs: Presentation, name: str):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    return None
