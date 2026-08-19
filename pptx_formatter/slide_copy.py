"""
Stage 3b - rebuilding a rough slide onto its matched layout.

Copying shapes across verbatim would preserve the rough deck's hand-placed
geometry and inherit nothing from the master, which defeats the point. So
content is *mapped into placeholders* instead: the slide's title text goes
into the layout's title placeholder, its body blocks into the body
placeholders in reading order, its picture into the picture placeholder.
Content that lands in a placeholder inherits the master's typography and
stays reflowable in PowerPoint, which is what makes the output genuinely
editable rather than a pile of pinned text boxes.

Not everything can or should be reflowed, so the mapping is a hybrid:

*   **Text, pictures** map into placeholders, keeping only inline emphasis
    (bold, italic, underline). Size, face and color are deliberately
    dropped so they resolve from the layout.
*   **Tables and charts** keep their own geometry but are moved onto the
    placeholder's footprint. A chart can't live in a text placeholder, and
    re-flowing a table would destroy its column widths.
*   **Anything left over** - decorative autoshapes, extra pictures - is
    admitted only if it can be placed without landing on mapped content.

That last rule is the one that keeps the output clean. Leftover shapes
arrive in the rough deck's coordinates while the placeholders sit on the
master's grid, so copying them verbatim is the one move guaranteed to
collide. Two gates decide their fate, and both fail closed - a shape that
can't be placed safely is dropped and reported, never nudged somewhere
arbitrary and hoped for:

1.  **Scaffolding is discarded.** A rough deck draws a card as a filled
    rectangle with separate text boxes on top of it. Once that text is
    mapped into a placeholder, the rectangle is an empty shell whose only
    job was framing words that now live somewhere else, and copying it
    across lands brand-less decoration on top of the layout that just
    received them. Whatever rode inside the shell goes with it.
2.  **Survivors are collision-checked.** Anything still standing is tested
    against every rectangle the mapped content now owns, and against the
    leftovers already admitted ahead of it.

Nothing is silently dropped: whatever can't be mapped is reported.
"""
from __future__ import annotations

import copy
import io

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.util import Emu

from . import geometry, part_copy
from .classifier import assign_roles

BODY_PLACEHOLDERS = {
    PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE,
}
TITLE_PLACEHOLDERS = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
LATENT_PLACEHOLDERS = {
    PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER,
}


# --- placeholder inventory ------------------------------------------------

def _sorted_by_position(shapes, rtl: bool):
    """Reading order: top to bottom, then left to right (or right to left)."""
    def key(shape):
        top = shape.top if shape.top is not None else 0
        left = shape.left if shape.left is not None else 0
        return (round(top / 100000), -left if rtl else left)
    return sorted(shapes, key=key)


def _inventory(slide, rtl: bool) -> dict:
    inv = {"title": None, "body": [], "picture": [], "table": [], "chart": []}
    for ph in slide.placeholders:
        ph_type = ph.placeholder_format.type
        if ph_type in TITLE_PLACEHOLDERS and inv["title"] is None:
            inv["title"] = ph
        elif ph_type in BODY_PLACEHOLDERS:
            inv["body"].append(ph)
        elif ph_type == PP_PLACEHOLDER.PICTURE:
            inv["picture"].append(ph)
        elif ph_type == PP_PLACEHOLDER.TABLE:
            inv["table"].append(ph)
        elif ph_type == PP_PLACEHOLDER.CHART:
            inv["chart"].append(ph)
    inv["body"] = _sorted_by_position(inv["body"], rtl)
    inv["picture"] = _sorted_by_position(inv["picture"], rtl)
    return inv


def _geometry(shape) -> tuple:
    return (shape.left, shape.top, shape.width, shape.height)


def _remove(shape) -> None:
    """Detach a shape, tolerating a shape that is already detached."""
    element = shape._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


# --- text mapping ---------------------------------------------------------

def _copy_paragraph(src_para, dest_para) -> None:
    """
    Copy one paragraph, keeping structure and inline emphasis but dropping
    typography.

    Face, size and color are intentionally not carried over: those are what
    the placeholder is supposed to supply. Bullet level, bold, italic and
    underline carry meaning rather than styling, so they survive.
    """
    if src_para.level:
        dest_para.level = src_para.level
    for src_run in src_para.runs:
        dest_run = dest_para.add_run()
        dest_run.text = src_run.text
        for attr in ("bold", "italic", "underline"):
            value = getattr(src_run.font, attr)
            if value is not None:
                setattr(dest_run.font, attr, value)


def copy_text(source, dest_frame) -> None:
    """Replace a placeholder's contents with a TextSource's paragraphs."""
    dest_frame.clear()
    for i, src_para in enumerate(source.para_list()):
        dest_para = dest_frame.paragraphs[0] if i == 0 else dest_frame.add_paragraph()
        _copy_paragraph(src_para, dest_para)


def _append_text(dest_frame, source) -> None:
    """Append a TextSource's paragraphs after whatever is already there."""
    for src_para in source.para_list():
        _copy_paragraph(src_para, dest_frame.add_paragraph())


def _shrink_on_overflow(text_frame) -> None:
    """
    Ask PowerPoint to scale text down when it overruns its slot.

    Only used where content was merged past what the layout budgeted for.
    Text set a size smaller than the designer intended is a compromise; text
    running off the bottom of the slide is a defect, and without this the
    merge produces the second one. PowerPoint recalculates the scale itself
    on open, so the element is written without a fontScale rather than
    guessing one here.
    """
    text_frame.word_wrap = True
    text_frame._txBody.bodyPr.get_or_change_to_normAutofit()


# --- leftover shapes ------------------------------------------------------

def _harvested_rects(title_source, body_sources) -> list:
    """
    Boxes of the source text boxes whose words are now in placeholders.

    A quote and its attribution can be split into two sources off one block,
    so the same shape may appear twice; duplicates are harmless here since
    these are only ever used as containment probes.
    """
    sources = list(body_sources)
    if title_source is not None:
        sources.append(title_source)

    rects = []
    for source in sources:
        block = getattr(source, "block", None)
        shape = getattr(block, "shape", None) if block is not None else None
        if shape is None:
            continue
        box = geometry.rect(shape)
        if box is not None:
            rects.append(box)
    return rects


def _scaffolding(leftovers: list, harvested: list) -> dict:
    """
    Leftover shapes that existed only to frame harvested text, keyed by id
    with the reason each one is going.

    Two passes. First the frames themselves: any leftover shape that
    contained one of the mapped text boxes. Then their contents: anything
    riding inside a doomed frame goes too, because a card's icon and divider
    rule floating on the layout without the card is worse than losing both.
    """
    frames = []
    for shape in leftovers:
        box = geometry.rect(shape)
        if box is None:
            continue
        if any(geometry.contains(box, text_box) for text_box in harvested):
            frames.append((shape, box))

    doomed = {
        id(shape): "framed text that is now in a placeholder"
        for shape, _ in frames
    }
    for shape in leftovers:
        if id(shape) in doomed:
            continue
        box = geometry.rect(shape)
        if box is None:
            continue
        if any(geometry.contains(frame_box, box) for _, frame_box in frames):
            doomed[id(shape)] = "rode inside a shape that was dropped"
    return doomed


# --- non-text mapping -----------------------------------------------------

def _place_at(shape, geometry) -> None:
    left, top, width, height = geometry
    if left is not None:
        shape.left = Emu(int(left))
    if top is not None:
        shape.top = Emu(int(top))
    if width is not None:
        shape.width = Emu(int(width))
    if height is not None:
        shape.height = Emu(int(height))


def _copy_table(dest_slide, source_shape, geometry=None):
    """
    Copy a table. Self-contained in XML, so a deep copy is sufficient - no
    external part to carry across.
    """
    new_el = copy.deepcopy(source_shape._element)
    dest_slide.shapes._spTree.append(new_el)
    shape = dest_slide.shapes[-1]
    if geometry and geometry[0] is not None:
        _place_at(shape, geometry)
    return shape


def _copy_chart(dest_slide, source_shape, geometry=None):
    """
    Copy a chart, cloning the chart part and its embedded workbook.

    Charts used to be skipped and replaced with a placeholder text box.
    They're carried properly now, since a consulting deck without its
    charts isn't a deck.
    """
    part_copy.clone_graphic_frame(dest_slide, source_shape)
    shape = dest_slide.shapes[-1]
    if geometry and geometry[0] is not None:
        _place_at(shape, geometry)
    return shape


def _copy_picture(dest_slide, source_shape, geometry=None):
    image = source_shape.image
    left, top, width, height = geometry if geometry else _geometry(source_shape)
    return dest_slide.shapes.add_picture(
        io.BytesIO(image.blob), Emu(int(left)), Emu(int(top)),
        Emu(int(width)), Emu(int(height)),
    )


def _copy_generic(dest_slide, source_shape):
    """Deep-copy a shape that references nothing outside its own XML."""
    new_el = copy.deepcopy(source_shape._element)
    dest_slide.shapes._spTree.append(new_el)
    return dest_slide.shapes[-1]


# --- page furniture -------------------------------------------------------

def clone_furniture(dest_slide, layout) -> list:
    """
    Copy footer, date and slide-number placeholders from the layout onto
    the slide.

    python-pptx treats these as "latent": they exist on the layout but are
    not cloned to new slides, which is why a deck can carry a perfectly good
    footer on every layout and still show nothing on any slide. PowerPoint
    clones them when the corresponding field is enabled, so we do too.
    """
    added = []
    existing = {
        ph.placeholder_format.type for ph in dest_slide.placeholders
    }
    for ph in layout.placeholders:
        ph_type = ph.placeholder_format.type
        if ph_type not in LATENT_PLACEHOLDERS or ph_type in existing:
            continue
        dest_slide.shapes._spTree.append(copy.deepcopy(ph._element))
        added.append(str(ph_type))
    return added


# --- entry point ----------------------------------------------------------

def rebuild_slide(dest_slide, source_slide, classification, spec, layout) -> dict:
    """
    Map a rough slide's content onto `dest_slide`, which was created from
    the layout matching `classification.archetype`.

    Returns a report of what was mapped, what was carried across
    unmapped, and anything that had to be improvised.
    """
    features = classification.features
    rtl_slide = features.is_rtl if features else False
    inv = _inventory(dest_slide, rtl_slide)
    report = {
        "mapped": [], "carried_over": [], "dropped": [],
        "warnings": [], "unused_placeholders": [],
    }

    title_source, body_sources = assign_roles(features, classification.archetype)
    if rtl_slide:
        # Right-to-left decks read right column first.
        body_sources.sort(key=lambda s: (round(s.block.top, 2), -s.block.left))

    # -- title -------------------------------------------------------------
    if title_source is not None and inv["title"] is not None:
        copy_text(title_source, inv["title"].text_frame)
        report["mapped"].append("title -> title placeholder")
    elif title_source is not None:
        # No title placeholder on this layout (a full-bleed picture, say);
        # keep the text rather than losing it.
        body_sources.insert(0, title_source)
        report["warnings"].append(
            f"layout '{layout.name}' has no title placeholder; title text kept as body content"
        )
    elif inv["title"] is not None:
        _remove(inv["title"])
        inv["title"] = None
        report["unused_placeholders"].append("title")

    # -- body text ---------------------------------------------------------
    body_slots = inv["body"]
    for i, slot in enumerate(body_slots):
        if i >= len(body_sources):
            break
        copy_text(body_sources[i], slot.text_frame)
        report["mapped"].append(f"body block {i + 1} -> {slot.name}")

    # More content than the layout has room for: fold the remainder into the
    # last slot rather than dropping it or spilling boxes over the design.
    if len(body_sources) > len(body_slots) and body_slots:
        for extra in body_sources[len(body_slots):]:
            _append_text(body_slots[-1].text_frame, extra)
        _shrink_on_overflow(body_slots[-1].text_frame)
        report["warnings"].append(
            f"{len(body_sources) - len(body_slots)} extra text block(s) merged into "
            f"'{body_slots[-1].name}'; layout '{layout.name}' has {len(body_slots)} body slot(s)"
        )
    elif body_sources and not body_slots:
        # No body slot at all - place the text in the grid's content area.
        left = int(spec.grid.margin_left_frac * spec.slide_width)
        top = int((spec.grid.margin_top_frac + 0.2) * spec.slide_height)
        width = int((1 - spec.grid.margin_left_frac - spec.grid.margin_right_frac) * spec.slide_width)
        height = int(0.5 * spec.slide_height)
        box = dest_slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        for i, source in enumerate(body_sources):
            if i == 0:
                copy_text(source, box.text_frame)
            else:
                _append_text(box.text_frame, source)
        _shrink_on_overflow(box.text_frame)
        report["warnings"].append(
            f"layout '{layout.name}' has no body placeholder; "
            f"{len(body_sources)} text block(s) placed in the content area"
        )

    # Body slots the text didn't need aren't discarded yet: a table or chart
    # on a layout with no dedicated slot of its own should be able to claim
    # one, which is what puts it on the design's footprint instead of
    # leaving it wherever the rough deck happened to draw it.
    spare_slots = list(body_slots[len(body_sources):])

    def claim(preferred: list) -> object | None:
        """Take the next available slot, preferring a type-specific one."""
        if preferred:
            return preferred.pop(0)
        return spare_slots.pop(0) if spare_slots else None

    # -- pictures ----------------------------------------------------------
    pictures = features.pictures if features else []
    pic_slots = list(inv["picture"])
    for i, picture in enumerate(pictures):
        slot = claim(pic_slots)
        if slot is None:
            _copy_picture(dest_slide, picture)
            report["carried_over"].append(
                f"picture '{picture.name}' kept at original position"
            )
            continue
        try:
            if slot.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                slot.insert_picture(io.BytesIO(picture.image.blob))
            else:
                # A text slot can't hold a picture, so use its footprint.
                footprint = _geometry(slot)
                _remove(slot)
                _copy_picture(dest_slide, picture, footprint)
            report["mapped"].append(f"picture {i + 1} -> {slot.name}")
        except (AttributeError, ValueError) as exc:
            report["warnings"].append(f"could not place picture in '{slot.name}': {exc}")

    # -- tables ------------------------------------------------------------
    tables = features.tables if features else []
    table_slots = list(inv["table"])
    for i, table in enumerate(tables):
        slot = claim(table_slots)
        footprint = None
        if slot is not None:
            footprint = _geometry(slot)
            _remove(slot)
        _copy_table(dest_slide, table, footprint)
        report["mapped"].append(
            f"table {i + 1} -> {'placeholder footprint' if footprint else 'original position'}"
        )

    # -- charts ------------------------------------------------------------
    charts = features.charts if features else []
    chart_slots = list(inv["chart"])
    for i, chart in enumerate(charts):
        slot = claim(chart_slots)
        footprint = None
        if slot is not None:
            footprint = _geometry(slot)
            _remove(slot)
        try:
            _copy_chart(dest_slide, chart, footprint)
            report["mapped"].append(
                f"chart {i + 1} copied with its data"
                + (" onto placeholder footprint" if footprint else "")
            )
        except Exception as exc:                      # noqa: BLE001
            # A malformed chart part shouldn't cost the whole deck.
            report["warnings"].append(f"could not copy chart '{chart.name}': {exc}")

    # Whatever no content claimed comes out, so the deck doesn't ship with
    # "Click to add text" prompts on it.
    for slot in spare_slots + pic_slots + table_slots + chart_slots:
        _remove(slot)
        report["unused_placeholders"].append(slot.name)

    # -- everything else ---------------------------------------------------
    # Both gates fail closed. A leftover shape earns its place on the slide
    # or it is dropped with a reason; it is never placed on the chance that
    # it happens to miss.
    leftovers = [
        shape
        for shape in (features.other_shapes if features else [])
        if shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER
    ]
    scaffolding = _scaffolding(leftovers, _harvested_rects(title_source, body_sources))

    # Everything already on the slide is mapped content, since leftovers are
    # copied last. Survivors join the set as they land, so they can't stack
    # on each other either.
    occupied = [box for box in (geometry.rect(s) for s in dest_slide.shapes) if box]

    for shape in leftovers:
        if id(shape) in scaffolding:
            report["dropped"].append(
                f"shape '{shape.name}' dropped: {scaffolding[id(shape)]}"
            )
            continue

        box = geometry.rect(shape)
        if box is not None and geometry.collides(box, occupied):
            report["dropped"].append(
                f"shape '{shape.name}' dropped: its position lands on mapped content"
            )
            continue

        try:
            _copy_generic(dest_slide, shape)
            report["carried_over"].append(f"shape '{shape.name}' copied at original position")
            if box is not None:
                occupied.append(box)
        except Exception as exc:                      # noqa: BLE001
            report["warnings"].append(f"could not copy shape '{shape.name}': {exc}")

    report["furniture_added"] = clone_furniture(dest_slide, layout)
    return report
