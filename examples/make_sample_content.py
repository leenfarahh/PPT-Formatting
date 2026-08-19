"""
Build a synthetic "rough content deck" for the demo.

Deliberately unstructured, the way a consultant's working deck arrives:
every slide is built on the blank layout out of hand-placed text boxes, so
nothing is in a placeholder and nothing inherits from a master. That is
exactly the input the classifier has to make sense of.

Covers the range the tool has to handle: a cover, bulleted content, a
two-column split, a comparison, a native chart, a table, an Arabic slide, a
pulled quote, a statistic, and a closing.

Run: python examples/make_sample_content.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

OUT_PATH = Path(__file__).resolve().parent / "sample_content.pptx"


def textbox(slide, text, left, top, width, height, size=None, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    lines = text.split("\n")
    frame.text = lines[0]
    for line in lines[1:]:
        frame.add_paragraph().text = line
    for para in frame.paragraphs:
        for run in para.runs:
            if size:
                run.font.size = Pt(size)
            run.font.bold = bold
    return box


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def new_slide():
        return prs.slides.add_slide(blank)

    # 1. Cover
    s = new_slide()
    textbox(s, "Q3 Growth Strategy", 1.0, 2.6, 9.0, 1.3, size=40, bold=True)
    textbox(s, "Prepared for the board  |  October 2026", 1.0, 4.1, 9.0, 0.6, size=16)

    # 2. Agenda
    s = new_slide()
    textbox(s, "Agenda", 0.8, 0.5, 10.0, 0.9, size=32, bold=True)
    textbox(
        s,
        "Market overview\nCompetitive position\nGrowth levers\nInvestment ask",
        0.8, 1.9, 10.0, 3.5, size=18,
    )

    # 3. Title and content
    s = new_slide()
    textbox(s, "Market Overview", 0.8, 0.5, 10.0, 0.9, size=32, bold=True)
    textbox(
        s,
        "Demand grew 14% year on year, ahead of the 9% forecast\n"
        "Competitor entry slowed materially through Q2\n"
        "Gross margins held above the 42% target in every region",
        0.8, 1.9, 10.0, 3.0, size=18,
    )

    # 4. Two column
    s = new_slide()
    textbox(s, "Where We Win", 0.8, 0.5, 11.0, 0.9, size=32, bold=True)
    textbox(s, "Enterprise accounts renew at 94% and expand at 1.3x net.",
            0.8, 2.1, 5.2, 2.6, size=18)
    textbox(s, "Mid-market is the growth engine but churn sits at 18%.",
            7.0, 2.1, 5.2, 2.6, size=18)

    # 5. Comparison
    s = new_slide()
    textbox(s, "Build or Partner", 0.8, 0.5, 11.0, 0.9, size=32, bold=True)
    textbox(s, "Build", 0.8, 1.9, 5.2, 0.5, size=20, bold=True)
    textbox(s, "Full control of roadmap, 14 months to parity, higher fixed cost.",
            0.8, 2.6, 5.2, 2.2, size=16)
    textbox(s, "Partner", 7.0, 1.9, 5.2, 0.5, size=20, bold=True)
    textbox(s, "Live in one quarter, revenue share of 22%, limited differentiation.",
            7.0, 2.6, 5.2, 2.2, size=16)

    # 6. Chart
    s = new_slide()
    textbox(s, "Revenue by Quarter", 0.8, 0.5, 10.0, 0.9, size=32, bold=True)
    data = CategoryChartData()
    data.categories = ["Q1", "Q2", "Q3", "Q4E"]
    data.add_series("Enterprise", (18.2, 22.5, 26.1, 29.4))
    data.add_series("Mid-market", (9.4, 10.1, 12.8, 15.2))
    s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.2), Inches(1.9), Inches(10.5), Inches(4.4), data,
    )

    # 7. Table
    s = new_slide()
    textbox(s, "Key Accounts", 0.8, 0.5, 10.0, 0.9, size=32, bold=True)
    rows = [
        ["Account", "Region", "Value (m)", "Status"],
        ["Northwind", "UAE", "1.24", "Renewed"],
        ["Contoso", "KSA", "0.96", "In review"],
        ["Fabrikam", "Qatar", "0.71", "At risk"],
    ]
    table = s.shapes.add_table(
        len(rows), len(rows[0]), Inches(1.2), Inches(2.0), Inches(10.5), Inches(2.6)
    ).table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            table.cell(r, c).text = value

    # 8. Arabic slide - exercises the complex-script and direction handling
    s = new_slide()
    textbox(s, "نظرة عامة على السوق", 0.8, 0.5, 10.0, 0.9, size=32, bold=True)
    textbox(
        s,
        "نمت المبيعات بنسبة 14% مقارنة بالعام الماضي\n"
        "تباطأ دخول المنافسين خلال الربع الثاني\n"
        "حافظنا على هوامش الربح فوق المستوى المستهدف",
        0.8, 1.9, 10.0, 3.0, size=18,
    )

    # 9. Quote
    s = new_slide()
    textbox(
        s,
        '"This partnership reset how our whole commercial team works."\n'
        "— Chief Executive, Northwind",
        1.6, 2.6, 9.5, 2.2, size=26,
    )

    # 10. Statistic
    s = new_slide()
    textbox(s, "94% net revenue retention", 1.6, 3.0, 10.0, 1.6, size=48, bold=True)

    # 11. Closing
    s = new_slide()
    textbox(s, "Thank you", 4.6, 3.2, 5.0, 1.4, size=44, bold=True)

    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH} ({len(prs.slides._sldIdLst)} rough slides)")


if __name__ == "__main__":
    main()
