"""
Build a synthetic "designer's master" for the demo.

Stands in for a real submission: a brand palette and typefaces written into
the theme, a logo on the master, footer text, layouts renamed the way a
designer would name them, and a couple of content slides that the tool is
expected to ignore.

Several standard layouts are deliberately deleted so the master has gaps -
that is what gives Stage 2 something to fill from the Template Bank.

Run: python examples/make_sample_master.py
"""
from __future__ import annotations

import struct
import sys
import zlib
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx_formatter.layout_builder import add_picture_to    # noqa: E402

OUT_PATH = Path(__file__).resolve().parent / "sample_master.pptx"
LOGO_PATH = Path(__file__).resolve().parent / "sample_logo.png"

A_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

# A plausible consulting-brand palette.
BRAND_COLORS = {
    "dk1": "1A1A2E", "lt1": "FFFFFF",
    "dk2": "16213E", "lt2": "F2F0EB",
    "accent1": "0F4C81",   # deep blue
    "accent2": "E8B04B",   # gold
    "accent3": "2E8B84",   # teal
    "accent4": "C1666B",   # clay
    "accent5": "6B7A8F",   # slate
    "accent6": "9BC4BC",   # sage
    "hlink": "0F4C81", "folHlink": "6B7A8F",
}
BRAND_FONTS = {
    "major_latin": "Georgia", "minor_latin": "Segoe UI",
    # Complex-script faces: what Arabic text resolves through.
    "major_cs": "Traditional Arabic", "minor_cs": "Dubai",
}

# Designer-style layout names, mapped onto the stock layouts they replace.
LAYOUT_RENAMES = {
    "Title Slide": "Cover",
    "Title and Content": "Standard Content",
    "Section Header": "Section Divider",
    "Two Content": "Two Column",
    "Comparison": "Comparison",
    "Title Only": "Title Only",
    "Picture with Caption": "Image with Text",
}
# Deleted so the submission has genuine gaps for Stage 2 to fill.
LAYOUTS_TO_DROP = [
    "Content with Caption", "Title and Vertical Text", "Vertical Title and Text",
]


def _png_chunk(tag: bytes, data: bytes) -> bytes:
    return (
        struct.pack(">I", len(data)) + tag + data
        + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
    )


def write_solid_png(path: Path, width: int, height: int, rgb: tuple) -> None:
    """
    Write a solid-color PNG without pulling in an image library.

    The demo only needs *a* logo file; keeping this dependency-free means
    `pip install -r requirements.txt` stays as small as it is.
    """
    raw = b"".join(b"\x00" + bytes(rgb) * width for _ in range(height))
    path.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        + _png_chunk(b"IHDR", struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0))
        + _png_chunk(b"IDAT", zlib.compress(raw, 9))
        + _png_chunk(b"IEND", b"")
    )


def apply_brand_theme(prs: Presentation) -> None:
    """Write the brand palette and typefaces into the theme part."""
    theme_part = prs.slide_masters[0].part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
    )
    root = etree.fromstring(theme_part.blob)

    scheme = root.find(".//a:clrScheme", A_NS)
    for role, hex_value in BRAND_COLORS.items():
        role_el = scheme.find(f"a:{role}", A_NS)
        if role_el is None:
            continue
        for child in list(role_el):
            role_el.remove(child)
        srgb = etree.SubElement(role_el, f"{{{A_NS['a']}}}srgbClr")
        srgb.set("val", hex_value)

    font_scheme = root.find(".//a:fontScheme", A_NS)
    for prefix, latin_key, cs_key in (
        ("major", "major_latin", "major_cs"),
        ("minor", "minor_latin", "minor_cs"),
    ):
        latin = font_scheme.find(f"a:{prefix}Font/a:latin", A_NS)
        cs = font_scheme.find(f"a:{prefix}Font/a:cs", A_NS)
        if latin is not None:
            latin.set("typeface", BRAND_FONTS[latin_key])
        if cs is not None:
            cs.set("typeface", BRAND_FONTS[cs_key])

    theme_part._blob = etree.tostring(
        root, xml_declaration=True, encoding="UTF-8", standalone=True
    )


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)   # 16:9
    apply_brand_theme(prs)

    master = prs.slide_masters[0]

    # Drop layouts so the submission has gaps to fill.
    for name in LAYOUTS_TO_DROP:
        layout = master.slide_layouts.get_by_name(name)
        if layout is not None:
            master.slide_layouts.remove(layout)

    for old, new in LAYOUT_RENAMES.items():
        layout = master.slide_layouts.get_by_name(old)
        if layout is not None:
            layout.name = new

    # A logo, small and in the top-right corner, which is where the
    # extractor's heuristic expects to find one.
    write_solid_png(LOGO_PATH, 240, 80, (15, 76, 129))
    add_picture_to(
        master.shapes, str(LOGO_PATH),
        Inches(11.4), Inches(0.35), Inches(1.5), Inches(0.5), name="Brand Logo",
    )

    for ph in master.placeholders:
        if "FOOTER" in str(ph.placeholder_format.type or "") and ph.has_text_frame:
            ph.text_frame.text = "Prezlab  |  Confidential"

    # Content slides that the tool must ignore.
    slide = prs.slides.add_slide(master.slide_layouts.get_by_name("Cover"))
    slide.shapes.title.text = "Brand guidelines deck (should be ignored)"
    slide = prs.slides.add_slide(master.slide_layouts.get_by_name("Standard Content"))
    slide.shapes.title.text = "Colour usage (should also be ignored)"

    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}")
    print(f"  logo:            {LOGO_PATH.name}")
    print(f"  layouts:         {[l.name for l in master.slide_layouts]}")
    print(f"  content slides:  {len(prs.slides._sldIdLst)} (to be ignored on ingest)")


if __name__ == "__main__":
    main()
