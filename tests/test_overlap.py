"""
Overlap prevention: leftover shapes must never land on mapped content.

A rough deck draws a card as a filled rectangle with separate text boxes on
top of it. Those coordinates mean nothing on the master's grid, so the
question every one of these tests asks is the same: given a rough slide that
would collide, does the output come out clean, and is what was removed
reported rather than silently vanished.

The overlap check here is written out longhand rather than imported from
pptx_formatter.geometry, so that a bug in the module under test cannot also
be the thing certifying its own output.
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches

from pptx_formatter.extraction import MARGIN_ASYMMETRY_LIMIT, derive_grid
from pptx_formatter.formatting import apply_grid_alignment
from pptx_formatter.pipeline import run_pipeline
from pptx_formatter.style_spec import LayoutSpec, PlaceholderSpec, StyleSpec

from conftest import textbox


# --- an independent overlap check ----------------------------------------

def _box(shape):
    if None in (shape.left, shape.top, shape.width, shape.height):
        return None
    return (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)


def _overlap_share(a, b):
    """Intersection as a share of the smaller box; 0.0 when they don't meet."""
    left, top = max(a[0], b[0]), max(a[1], b[1])
    right, bottom = min(a[2], b[2]), min(a[3], b[3])
    if right <= left or bottom <= top:
        return 0.0
    overlap = (right - left) * (bottom - top)
    areas = [(x[2] - x[0]) * (x[3] - x[1]) for x in (a, b)]
    smaller = min(areas)
    return 0.0 if smaller <= 0 else overlap / smaller


def overlapping_pairs(slide, tolerance=0.10):
    """Every pair of shapes covering more than `tolerance` of the smaller."""
    boxes = [(s, _box(s)) for s in slide.shapes]
    boxes = [(s, b) for s, b in boxes if b and (b[2] - b[0]) > 0 and (b[3] - b[1]) > 0]
    found = []
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            share = _overlap_share(boxes[i][1], boxes[j][1])
            if share > tolerance:
                found.append((boxes[i][0].name, boxes[j][0].name, share))
    return found


# --- a rough deck built the way the broken one was ------------------------

def _decoration(slide, shape_type, left, top, width, height, rgb=(0xE6, 0xC8, 0x96)):
    """A filled shape carrying no text of its own: pure decoration."""
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)
    shape.line.fill.background()
    shape.text_frame.clear()
    return shape


@pytest.fixture
def card_content_path(tmp_path):
    """
    A rough deck of card slides: a tan rounded rectangle per column with an
    icon circle and a divider rule inside it, and the words in separate text
    boxes laid on top. This is the shape that produced the overlapping
    output, reduced to its essentials.
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    textbox(slide, "Full Lifecycle Management", 0.85, 0.5, 9.0, 0.9, size=32, bold=True)
    textbox(slide, "CORE FEATURES", 0.85, 1.45, 4.0, 0.35, size=12, bold=True)
    textbox(slide, "Every entity supports the same operations.", 0.85, 1.85, 11.6, 0.7, size=14)

    for i, (heading, body) in enumerate([
        ("Workstreams", "Create, rename, archive and restore."),
        ("Projects", "Group tasks, set owners, track status."),
        ("Tasks", "Assign, comment, attach, move."),
    ]):
        left = 0.85 + i * 4.0
        _decoration(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, 2.8, 3.6, 3.6)
        _decoration(slide, MSO_SHAPE.OVAL, left + 0.3, 3.05, 0.65, 0.65,
                    rgb=(0x11, 0x87, 0x8B))
        _decoration(slide, MSO_SHAPE.RECTANGLE, left + 0.3, 4.0, 3.0, 0.01,
                    rgb=(0xFF, 0xFF, 0xFF))
        textbox(slide, heading, left + 0.3, 4.15, 3.0, 0.4, size=16, bold=True)
        textbox(slide, body, left + 0.3, 4.6, 3.0, 1.6, size=12)

    path = tmp_path / "cards.pptx"
    prs.save(str(path))
    return path


# --- the end-to-end guarantee --------------------------------------------

def test_card_deck_comes_out_with_nothing_overlapping(
    master_path, card_content_path, bank, tmp_path
):
    out = tmp_path / "out.pptx"
    run_pipeline(master_path, card_content_path, out, bank=bank, client="Cards")

    prs = Presentation(str(out))
    for index, slide in enumerate(prs.slides, 1):
        pairs = overlapping_pairs(slide)
        assert not pairs, f"slide {index} has overlapping shapes: {pairs}"


def test_the_emptied_cards_are_reported_not_silently_removed(
    master_path, card_content_path, bank, tmp_path
):
    """A shape that leaves the deck has to say so and say why."""
    report = run_pipeline(
        master_path, card_content_path, tmp_path / "out.pptx",
        bank=bank, client="Cards",
    )

    dropped = [line for slide in report["slides"] for line in slide["dropped"]]
    assert dropped, "the cards framed mapped text and should have been dropped"
    assert any("framed text" in line for line in dropped)
    # The icon and the rule rode inside a card, so they go with it.
    assert any("rode inside" in line for line in dropped)
    assert all("dropped:" in line for line in dropped), dropped


# --- scaffolding versus genuine decoration -------------------------------

def test_scaffolding_goes_even_when_it_would_not_collide(widescreen):
    """
    The two gates are not redundant.

    On a card-heavy slide the collision gate happens to catch the emptied
    cards as well, which can mask whether the scaffolding gate works at all.
    Here the card sits clear of every placeholder on the destination layout,
    so nothing but the scaffolding gate can remove it - and it still has to,
    because an empty shell whose words now live in a placeholder is not
    decoration, it is litter.
    """
    from pptx_formatter.classifier import classify_slide
    from pptx_formatter.slide_copy import rebuild_slide

    source = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(source, "Market Overview", 0.8, 0.5, 9.0, 0.9, size=32, bold=True)
    textbox(source, "Demand grew 14 percent.", 0.8, 1.9, 8.0, 1.2, size=18)

    # A card in the bottom-right corner, clear of the layout's placeholders,
    # with its words in a text box laid inside it.
    card = _decoration(source, MSO_SHAPE.ROUNDED_RECTANGLE, 10.4, 6.6, 2.6, 0.7)
    textbox(source, "Footnote copy", 10.5, 6.7, 2.4, 0.5, size=11)

    classification = classify_slide(
        source, 0, 1, widescreen.slide_width, widescreen.slide_height
    )
    dest = widescreen.slides.add_slide(widescreen.slide_layouts[1])
    spec = StyleSpec()
    spec.slide_width, spec.slide_height = widescreen.slide_width, widescreen.slide_height
    spec.grid.compute_guides()

    # The card really is clear of the destination placeholders, or this test
    # would be re-testing the collision gate by accident.
    card_box = (card.left, card.top, card.left + card.width, card.top + card.height)
    for ph in dest.placeholders:
        ph_box = _box(ph)
        if ph_box:
            assert _overlap_share(card_box, ph_box) == 0.0, "card is not clear"

    report = rebuild_slide(
        dest, source, classification, spec, widescreen.slide_layouts[1]
    )

    reasons = [line for line in report["dropped"] if card.name in line]
    assert reasons, f"the emptied card survived: {report['dropped']}"
    assert "framed text" in reasons[0], reasons



def test_decoration_clear_of_mapped_content_survives(widescreen, tmp_path):
    """
    The guard drops what collides, not everything it doesn't understand.

    A shape sitting in whitespace that frames no mapped text has no reason to
    go, and dropping it would quietly strip brand furniture off the deck.
    """
    from pptx_formatter.classifier import classify_slide, extract_features
    from pptx_formatter.slide_copy import rebuild_slide

    source = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(source, "Market Overview", 0.8, 0.5, 10, 0.9, size=32, bold=True)
    textbox(source, "Demand grew 14 percent.", 0.8, 1.9, 10, 2.0, size=18)
    # A rule in the bottom margin: touches nothing that gets mapped.
    keeper = _decoration(source, MSO_SHAPE.RECTANGLE, 0.8, 7.1, 2.0, 0.05)

    classification = classify_slide(
        source, 0, 1, widescreen.slide_width, widescreen.slide_height
    )
    dest = widescreen.slides.add_slide(widescreen.slide_layouts[1])
    spec = StyleSpec()
    spec.slide_width, spec.slide_height = widescreen.slide_width, widescreen.slide_height
    spec.grid.compute_guides()

    report = rebuild_slide(
        dest, source, classification, spec, widescreen.slide_layouts[1]
    )

    assert not any(keeper.name in line for line in report["dropped"])
    assert any(keeper.name in line for line in report["carried_over"])


# --- the grid pass may not undo the work ---------------------------------

def test_snapping_is_abandoned_when_it_would_create_an_overlap(widescreen):
    """
    A snap that pushes a shape onto a placeholder is refused outright.

    Off the guide by a few millimetres is a smaller defect than content
    buried under content, so the guard fails closed rather than hunting for
    somewhere else to put the shape.
    """
    spec = StyleSpec()
    spec.slide_width, spec.slide_height = widescreen.slide_width, widescreen.slide_height
    spec.grid.compute_guides()

    slide = widescreen.slides.add_slide(widescreen.slide_layouts[1])
    body = slide.placeholders[1]
    # Park a shape just off the guide, squarely over the body placeholder.
    intruder = textbox(
        slide, "Intruder",
        Emu(body.left).inches + 0.07, Emu(body.top).inches + 0.07,
        2.0, 1.0, size=14,
    )
    before = (intruder.left, intruder.top)

    moved = apply_grid_alignment(slide, spec)

    assert intruder.name not in moved
    assert (intruder.left, intruder.top) == before


def test_snapping_still_aligns_a_shape_with_room_to_move(widescreen):
    """The guard must not turn the grid pass into a no-op."""
    spec = StyleSpec()
    spec.slide_width, spec.slide_height = widescreen.slide_width, widescreen.slide_height
    spec.grid.compute_guides()

    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    box = textbox(slide, "Adrift", 1.07, 1.03, 3.0, 0.8, size=18)
    before = (box.left, box.top)

    moved = apply_grid_alignment(slide, spec)

    assert box.name in moved
    assert (box.left, box.top) != before


# --- grid inference ------------------------------------------------------

def _phs(*specs):
    return [
        PlaceholderSpec(ph_type=t, idx=i, left_frac=l, top_frac=0.3,
                        width_frac=w, height_frac=0.4)
        for i, (t, l, w) in enumerate(specs)
    ]


def test_right_margin_is_mirrored_when_nothing_evidences_it():
    """
    No full-width placeholder means no evidence for the right margin.

    Deriving it from the widest thing that happened to be measured strands a
    band of dead slide down the side of every generated layout, which is what
    a master carrying its design in plain shapes produces.
    """
    layouts = [LayoutSpec(placeholders=_phs(
        ("title", 0.064, 0.55),
        ("body", 0.29, 0.41),
        ("pic", 0.44, 0.17),
    ))]

    grid = derive_grid(layouts)

    # Widest right edge is 0.70, so the raw inference wants 0.30.
    assert grid.margin_right_frac == pytest.approx(grid.margin_left_frac)
    usable = 1.0 - grid.margin_left_frac - grid.margin_right_frac
    assert usable > 0.85, f"content box collapsed to {usable:.0%} of the slide"


def test_a_genuine_asymmetric_grid_is_left_alone():
    """Within the limit, the designer's own margins stand."""
    layouts = [LayoutSpec(placeholders=_phs(
        ("title", 0.1, 0.8),
        ("body", 0.1, 0.38),
        ("body", 0.52, 0.38),
    ))]

    grid = derive_grid(layouts)

    assert grid.margin_left_frac == pytest.approx(0.1)
    assert grid.margin_right_frac == pytest.approx(0.1)
    assert grid.margin_right_frac <= grid.margin_left_frac * MARGIN_ASYMMETRY_LIMIT


def test_too_few_placeholders_falls_back_to_defaults():
    """Two stray frames describe their own layout, not the deck."""
    defaults = StyleSpec().grid
    layouts = [LayoutSpec(placeholders=_phs(("pic", 0.44, 0.17), ("pic", 0.60, 0.18)))]

    grid = derive_grid(layouts)

    assert grid.margin_left_frac == defaults.margin_left_frac
    assert grid.margin_right_frac == defaults.margin_right_frac
