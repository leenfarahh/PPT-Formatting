"""Stage 3c: typography, color, direction and grid handling."""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from pptx_formatter import rtl
from pptx_formatter.formatting import (
    MINOR_CS, MINOR_LATIN, apply_color_mapping, apply_grid_alignment,
    apply_typography, format_table, nearest_accent_role,
)
from pptx_formatter.style_spec import StyleSpec

from conftest import textbox

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def brand_spec():
    spec = StyleSpec()
    spec.theme.colors.update({
        "dk1": "1A1A2E", "lt1": "FFFFFF",
        "accent1": "0F4C81", "accent2": "E8B04B", "accent3": "2E8B84",
    })
    spec.slide_width, spec.slide_height = Inches(13.333), Inches(7.5)
    spec.grid.compute_guides()
    return spec


# --- Arabic and direction -------------------------------------------------

def test_rtl_detection():
    assert rtl.is_rtl_text("نظرة عامة على السوق")
    assert not rtl.is_rtl_text("Market Overview")
    # Neutral characters shouldn't dilute the signal.
    assert rtl.is_rtl_text("نمت المبيعات بنسبة 14% (2026)")
    assert rtl.rtl_ratio("") == 0.0


def test_arabic_sentence_quoting_english_stays_rtl():
    """An Arabic sentence naming an English product is still Arabic."""
    assert rtl.is_rtl_text("أطلقنا منتج Growth Suite في السوق")


def test_paragraph_direction_and_alignment_applied(widescreen):
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, "نمت المبيعات بنسبة كبيرة", 1, 1, 8, 1, size=18)

    stats = apply_typography(slide, brand_spec())
    para = slide.shapes[0].text_frame.paragraphs[0]
    pPr = para._p.find(qn("a:pPr"))

    assert stats["rtl_paragraphs"] == 1
    assert pPr.get("rtl") == "1"
    assert pPr.get("algn") == "r"


def test_latin_paragraph_is_not_flipped(widescreen):
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, "Market Overview", 1, 1, 8, 1, size=18)
    apply_typography(slide, brand_spec())
    pPr = slide.shapes[0].text_frame.paragraphs[0]._p.find(qn("a:pPr"))
    assert pPr.get("rtl") == "0"


def test_centered_text_is_left_alone(widescreen):
    """Centering is a design decision, not a direction default."""
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    box = textbox(slide, "نظرة عامة", 1, 1, 8, 1, size=18)
    from pptx.enum.text import PP_ALIGN
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    apply_typography(slide, brand_spec())
    assert box.text_frame.paragraphs[0]._p.find(qn("a:pPr")).get("algn") == "ctr"


def test_both_scripts_get_a_typeface(widescreen):
    """
    A bilingual run needs the Latin face and the Arabic face set, because
    PowerPoint resolves the font per character.
    """
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, "أطلقنا Growth Suite", 1, 1, 8, 1, size=18)
    apply_typography(slide, brand_spec())

    rPr = slide.shapes[0].text_frame.paragraphs[0].runs[0].font._rPr
    assert rPr.find(qn("a:latin")).get("typeface") == MINOR_LATIN
    assert rPr.find(qn("a:cs")).get("typeface") == MINOR_CS


def test_theme_references_used_rather_than_literal_fonts(widescreen):
    """Pinning a literal face would undo the inheritance Stage 2 sets up."""
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, "Market Overview", 1, 1, 8, 1, size=18)
    apply_typography(slide, brand_spec())

    typeface = slide.shapes[0].text_frame.paragraphs[0].runs[0].font._rPr.find(
        qn("a:latin")
    ).get("typeface")
    assert typeface.startswith("+"), "font should reference the theme, not a literal name"


# --- color ----------------------------------------------------------------

def test_nearest_accent_matching():
    spec = brand_spec()
    assert nearest_accent_role("0F4C82", spec) == "accent1"    # a shade off the blue
    assert nearest_accent_role("E8B04A", spec) == "accent2"


def test_stray_fills_remapped_to_the_palette(widescreen):
    from pptx.enum.shapes import MSO_SHAPE
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string("0F4C90")   # near accent1

    notes = apply_color_mapping(slide, brand_spec())
    assert str(shape.fill.fore_color.rgb) == "0F4C81"
    assert len(notes) == 1


def test_placeholder_fills_are_not_touched(master_path):
    """A placeholder's fill comes from the layout; overwriting it pins it."""
    from pptx_formatter.extraction import extract_style_spec
    from pptx_formatter.layout_generator import generate_master_layouts

    spec = extract_style_spec(master_path)
    prs, _ = generate_master_layouts(spec, master_path)
    slide = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[1])
    slide.shapes.title.text = "Title"

    assert apply_color_mapping(slide, spec) == []


# --- grid -----------------------------------------------------------------

def test_unmapped_shapes_snap_to_the_grid(widescreen):
    spec = brand_spec()
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    box = textbox(slide, "Adrift", 1.07, 1.03, 4, 1, size=18)
    before = (box.left, box.top)

    moved = apply_grid_alignment(slide, spec)
    assert box.name in moved
    assert (box.left, box.top) != before
    # Snapped, but still inside the margins.
    assert box.left >= spec.grid.margin_left_frac * spec.slide_width - 1


def test_placeholders_are_never_moved(master_path):
    """They already sit exactly where the layout puts them."""
    from pptx_formatter.extraction import extract_style_spec
    from pptx_formatter.layout_generator import generate_master_layouts

    spec = extract_style_spec(master_path)
    prs, _ = generate_master_layouts(spec, master_path)
    slide = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[1])
    slide.shapes.title.text = "Title"
    before = [(ph.left, ph.top) for ph in slide.placeholders]

    apply_grid_alignment(slide, spec)
    assert [(ph.left, ph.top) for ph in slide.placeholders] == before


def test_snapping_keeps_shapes_on_the_slide(widescreen):
    spec = brand_spec()
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    box = textbox(slide, "Far right", 12.9, 7.2, 3, 1, size=18)

    apply_grid_alignment(slide, spec)
    assert box.left + box.width <= spec.slide_width
    assert box.top + box.height <= spec.slide_height


# --- tables ---------------------------------------------------------------

def test_table_styled_from_the_spec(widescreen):
    spec = brand_spec()
    spec.table_style.header_fill = "0F4C81"
    spec.table_style.header_font_color = "FFFFFF"
    spec.table_style.banded_fill = None

    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    shape = slide.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(8), Inches(2))
    table = shape.table
    for r in range(3):
        for c in range(2):
            table.cell(r, c).text = f"r{r}c{c}"

    format_table(shape, spec)
    header_run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
    body_run = table.cell(1, 0).text_frame.paragraphs[0].runs[0]

    assert str(table.cell(0, 0).fill.fore_color.rgb) == "0F4C81"
    assert header_run.font.bold
    assert header_run.font.size == Pt(spec.table_style.header_font_size_pt)
    assert body_run.font.size == Pt(spec.table_style.body_font_size_pt)


def test_arabic_table_cells_get_direction(widescreen):
    spec = brand_spec()
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    shape = slide.shapes.add_table(2, 1, Inches(1), Inches(1), Inches(8), Inches(2))
    shape.table.cell(0, 0).text = "المنطقة"
    shape.table.cell(1, 0).text = "الإمارات"

    format_table(shape, spec)
    pPr = shape.table.cell(1, 0).text_frame.paragraphs[0]._p.find(qn("a:pPr"))
    assert pPr.get("rtl") == "1"
