"""
Shared fixtures and the package validator.

`assert_valid_pptx` is the important one. There's no PowerPoint (or
LibreOffice) in CI to open the output, and the failure mode that matters
most for this tool is a structurally broken package - a shape referencing a
relationship that doesn't exist, or a part missing from `[Content_Types]`.
Both produce a file that python-pptx reopens quite happily and PowerPoint
refuses with "needs to be repaired". Checking the package directly is what
catches that class of bug.
"""
from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"

# Attributes that name a relationship from within part XML.
REL_ATTRS = [f"{{{R_NS}}}id", f"{{{R_NS}}}embed", f"{{{R_NS}}}link"]


def _rels_name(part_name: str) -> str:
    path = Path(part_name)
    return f"{path.parent.as_posix()}/_rels/{path.name}.rels"


def assert_valid_pptx(path) -> None:
    """
    Check a .pptx is structurally sound.

    Verifies three things across every XML part: it parses, every
    relationship it references by rId exists in that part's `.rels`, and
    every part is covered by `[Content_Types].xml`.
    """
    path = str(path)
    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
        assert "[Content_Types].xml" in names, "package has no [Content_Types].xml"

        content_types = etree.fromstring(archive.read("[Content_Types].xml"))
        defaults = {
            el.get("Extension").lower()
            for el in content_types.findall(f"{{{CT_NS}}}Default")
        }
        overrides = {
            el.get("PartName") for el in content_types.findall(f"{{{CT_NS}}}Override")
        }

        for name in names:
            if name == "[Content_Types].xml" or name.endswith("/"):
                continue

            # Every part must be typed, by extension default or by override.
            part_name = "/" + name
            extension = name.rsplit(".", 1)[-1].lower() if "." in name else ""
            assert part_name in overrides or extension in defaults, (
                f"{name} has no content type declared"
            )

            if not name.endswith((".xml", ".rels")):
                continue

            try:
                root = etree.fromstring(archive.read(name))
            except etree.XMLSyntaxError as exc:                # pragma: no cover
                raise AssertionError(f"{name} is not well-formed XML: {exc}") from exc

            if name.endswith(".rels"):
                continue

            referenced = {
                el.get(attr)
                for el in root.iter()
                for attr in REL_ATTRS
                if el.get(attr)
            }
            if not referenced:
                continue

            rels_name = _rels_name(name)
            assert rels_name in names, (
                f"{name} references {sorted(referenced)} but has no .rels part"
            )
            rels_root = etree.fromstring(archive.read(rels_name))
            declared = {
                el.get("Id") for el in rels_root.findall(f"{{{REL_NS}}}Relationship")
            }
            missing = referenced - declared
            assert not missing, f"{name} references undeclared relationships: {sorted(missing)}"


def textbox(slide, text, left, top, width, height, size=None, bold=False):
    """Add a plain text box, the way a rough deck is actually built."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    lines = text.split("\n")
    frame.text = lines[0]
    for line in lines[1:]:
        frame.add_paragraph().text = line
    for para in frame.paragraphs:
        for run in para.runs:
            if size:
                run.font.size = Pt(size)
            if bold:
                run.font.bold = True
    return box


@pytest.fixture
def widescreen():
    """A blank 16:9 presentation."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    return prs


@pytest.fixture
def master_path(tmp_path):
    """
    A stand-in for a designer's master: brand colors and fonts in the
    theme, two layouts removed to leave gaps, and content slides present
    that ingestion is expected to ignore.
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    master = prs.slide_masters[0]
    theme_part = master.part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
    )
    root = etree.fromstring(theme_part.blob)
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"

    colors = {
        "dk1": "1A1A2E", "lt1": "FFFFFF", "dk2": "16213E", "lt2": "F2F0EB",
        "accent1": "0F4C81", "accent2": "E8B04B", "accent3": "2E8B84",
        "accent4": "C1666B", "accent5": "6B7A8F", "accent6": "9BC4BC",
    }
    scheme = root.find(f".//{{{a}}}clrScheme")
    for role, value in colors.items():
        role_el = scheme.find(f"{{{a}}}{role}")
        for child in list(role_el):
            role_el.remove(child)
        etree.SubElement(role_el, f"{{{a}}}srgbClr").set("val", value)

    font_scheme = root.find(f".//{{{a}}}fontScheme")
    for prefix, latin, cs in (("major", "Georgia", "Traditional Arabic"),
                              ("minor", "Segoe UI", "Dubai")):
        font_scheme.find(f"{{{a}}}{prefix}Font/{{{a}}}latin").set("typeface", latin)
        font_scheme.find(f"{{{a}}}{prefix}Font/{{{a}}}cs").set("typeface", cs)

    theme_part._blob = etree.tostring(
        root, xml_declaration=True, encoding="UTF-8", standalone=True
    )

    for name in ("Content with Caption", "Title and Vertical Text"):
        layout = master.slide_layouts.get_by_name(name)
        if layout is not None:
            master.slide_layouts.remove(layout)

    for ph in master.placeholders:
        if "FOOTER" in str(ph.placeholder_format.type or "") and ph.has_text_frame:
            ph.text_frame.text = "Prezlab | Confidential"

    # Content slides that must be ignored on ingest.
    slide = prs.slides.add_slide(master.slide_layouts[0])
    slide.shapes.title.text = "Ignore me"
    prs.slides.add_slide(master.slide_layouts[1])

    path = tmp_path / "master.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def content_path(tmp_path):
    """A rough deck: hand-placed text boxes, nothing in a placeholder."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    s = prs.slides.add_slide(blank)                       # 1 cover
    textbox(s, "Q3 Growth Strategy", 1, 2.6, 9, 1.3, size=40, bold=True)
    textbox(s, "Prepared for the board", 1, 4.1, 9, 0.6, size=16)

    s = prs.slides.add_slide(blank)                       # 2 title + content
    textbox(s, "Market Overview", 0.8, 0.5, 10, 0.9, size=32, bold=True)
    textbox(s, "Demand grew 14%\nCompetitors slowed\nMargins held", 0.8, 1.9, 10, 3, size=18)

    s = prs.slides.add_slide(blank)                       # 3 two column
    textbox(s, "Where We Win", 0.8, 0.5, 11, 0.9, size=32, bold=True)
    textbox(s, "Enterprise renews at 94% and expands.", 0.8, 2.1, 5.2, 2.6, size=18)
    textbox(s, "Mid-market grows fast but churns.", 7.0, 2.1, 5.2, 2.6, size=18)

    s = prs.slides.add_slide(blank)                       # 4 arabic
    textbox(s, "نظرة عامة على السوق", 0.8, 0.5, 10, 0.9, size=32, bold=True)
    textbox(s, "نمت المبيعات بنسبة 14% مقارنة بالعام الماضي", 0.8, 1.9, 10, 3, size=18)

    s = prs.slides.add_slide(blank)                       # 5 quote
    textbox(s, '"This reset how we work."\n— Chief Executive', 1.6, 2.6, 9.5, 2.2, size=26)

    s = prs.slides.add_slide(blank)                       # 6 closing
    textbox(s, "Thank you", 4.6, 3.2, 5, 1.4, size=44, bold=True)

    path = tmp_path / "content.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def bank(tmp_path):
    from pptx_formatter.bank import TemplateBank
    return TemplateBank(tmp_path / "bank")
