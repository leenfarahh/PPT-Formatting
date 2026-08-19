"""Stage 1: reading a master into a Style Spec."""
from __future__ import annotations

import json

from pptx import Presentation

from pptx_formatter import archetypes
from pptx_formatter.extraction import (
    classify_layout, derive_grid, extract_style_spec,
)
from pptx_formatter.style_spec import PlaceholderSpec, SPEC_VERSION, StyleSpec


def test_content_slides_are_ignored(master_path):
    """A master may arrive full; only its layouts matter."""
    spec = extract_style_spec(master_path)
    assert spec.meta.content_slides_ignored == 2
    assert spec.meta.layouts_found == 9
    assert len(spec.layouts) == 9


def test_extracts_theme_colors_and_per_script_fonts(master_path):
    spec = extract_style_spec(master_path)

    assert spec.theme.colors["accent1"] == "0F4C81"
    assert spec.theme.colors["dk1"] == "1A1A2E"
    assert spec.theme.fonts.major_latin == "Georgia"
    assert spec.theme.fonts.minor_latin == "Segoe UI"
    # Arabic resolves through the complex-script face, so it must survive.
    assert spec.theme.fonts.minor_cs == "Dubai"
    assert spec.theme.fonts.cs_for(is_major=True) == "Traditional Arabic"


def test_complex_script_falls_back_to_latin_when_undeclared():
    fonts = StyleSpec().theme.fonts
    fonts.major_latin = "Helvetica"
    assert fonts.major_cs == ""
    assert fonts.cs_for(is_major=True) == "Helvetica"


def test_placeholder_geometry_is_normalized_to_fractions(master_path):
    spec = extract_style_spec(master_path)
    for layout in spec.layouts:
        for ph in layout.placeholders:
            assert 0.0 <= ph.left_frac <= 1.0
            assert 0.0 <= ph.top_frac <= 1.0
            assert 0.0 < ph.width_frac <= 1.0


def test_fractions_resolve_against_a_different_slide_size():
    """A layout banked from 4:3 has to land correctly on 16:9."""
    ph = PlaceholderSpec(left_frac=0.5, top_frac=0.25, width_frac=0.25, height_frac=0.5)
    left, top, width, height = ph.to_emu(12192000, 6858000)
    assert left == 6096000
    assert top == 1714500
    assert width == 3048000
    assert height == 3429000


def test_layouts_are_tagged_with_archetypes(master_path):
    spec = extract_style_spec(master_path)
    present = spec.archetypes_present()
    assert archetypes.TITLE_SLIDE in present
    assert archetypes.TITLE_AND_CONTENT in present
    assert archetypes.TWO_CONTENT in present
    # Deleting layouts from the master must leave real gaps to fill.
    assert archetypes.QUOTE not in present
    assert archetypes.CHART not in present


def test_designer_layout_name_beats_structure():
    """A layout the designer called "Quote" is a quote layout."""
    placeholders = [
        PlaceholderSpec(ph_type="title", width_frac=0.8, height_frac=0.2),
        PlaceholderSpec(ph_type="body", idx=1, width_frac=0.8, height_frac=0.5),
    ]
    assert classify_layout("Pull Quote", "obj", placeholders) == archetypes.QUOTE
    assert classify_layout("Section Divider", "obj", placeholders) == archetypes.SECTION_HEADER
    assert classify_layout("Standard Content", "obj", placeholders) == \
        archetypes.TITLE_AND_CONTENT


def test_full_bleed_picture_detected_from_geometry():
    placeholders = [PlaceholderSpec(ph_type="pic", idx=1, width_frac=1.0, height_frac=1.0)]
    assert classify_layout("Layout 7", "obj", placeholders) == archetypes.PICTURE_FULL


def test_grid_is_derived_from_placeholder_geometry():
    """Margins come from where the designer actually put things."""
    from pptx_formatter.style_spec import LayoutSpec

    layouts = [LayoutSpec(placeholders=[
        PlaceholderSpec(ph_type="title", left_frac=0.1, top_frac=0.08,
                        width_frac=0.8, height_frac=0.15),
        PlaceholderSpec(ph_type="body", idx=1, left_frac=0.1, top_frac=0.3,
                        width_frac=0.38, height_frac=0.5),
        PlaceholderSpec(ph_type="body", idx=2, left_frac=0.52, top_frac=0.3,
                        width_frac=0.38, height_frac=0.5),
    ])]
    grid = derive_grid(layouts)

    assert grid.margin_left_frac == 0.1
    assert grid.margin_top_frac == 0.08
    assert abs(grid.margin_right_frac - 0.1) < 1e-6
    # The gap between the two side-by-side bodies is the gutter.
    assert abs(grid.gutter_frac - 0.04) < 1e-6
    assert grid.column_guides and grid.column_guides[0] == 0.1


def test_chart_table_and_icon_styles_derive_from_the_theme(master_path):
    spec = extract_style_spec(master_path)

    assert spec.chart_style.series_colors[0] == "0F4C81"
    assert len(spec.chart_style.series_colors) == 6
    assert spec.table_style.header_fill == "0F4C81"
    # A dark header needs light text on it.
    assert spec.table_style.header_font_color == "FFFFFF"
    assert spec.icon_palette == spec.accent_hex_list()


def test_meta_records_provenance(master_path):
    spec = extract_style_spec(master_path, client="Acme", project="Board Deck")
    assert spec.meta.spec_version == SPEC_VERSION
    assert spec.meta.client == "Acme"
    assert spec.meta.project == "Board Deck"
    assert spec.meta.extracted_at is not None
    assert spec.meta.source_name == "master"


def test_logo_is_extracted_with_its_bytes(tmp_path, widescreen):
    """Geometry alone isn't enough; Stage 2 needs the image itself."""
    from pptx.util import Inches
    from pptx_formatter.layout_builder import add_picture_to
    from examples.make_sample_master import write_solid_png

    logo_src = tmp_path / "logo.png"
    write_solid_png(logo_src, 120, 40, (15, 76, 129))
    add_picture_to(
        widescreen.slide_masters[0].shapes, str(logo_src),
        Inches(11.4), Inches(0.35), Inches(1.5), Inches(0.5), name="Brand Logo",
    )
    master = tmp_path / "with_logo.pptx"
    widescreen.save(str(master))

    spec = extract_style_spec(master, asset_dir=tmp_path / "assets")
    assert spec.brand.logo.present
    assert spec.brand.logo.asset_path.endswith(".png")
    assert 0.8 < spec.brand.logo.left_frac < 0.9


def test_spec_round_trips_through_json(master_path, tmp_path):
    """The spec is the durable artifact, so it has to survive a round trip."""
    spec = extract_style_spec(master_path, client="Acme")
    path = tmp_path / "spec.json"
    spec.save(path)

    reloaded = StyleSpec.load(path)
    assert reloaded.theme.colors == spec.theme.colors
    assert reloaded.theme.fonts.minor_cs == spec.theme.fonts.minor_cs
    assert reloaded.meta.client == "Acme"
    assert len(reloaded.layouts) == len(spec.layouts)
    assert reloaded.layouts[0].placeholders[0].ph_type == spec.layouts[0].placeholders[0].ph_type


def test_json_shape_matches_the_documented_contract(master_path):
    """The Style Spec is consumed by other tools, so its shape is a contract."""
    document = json.loads(extract_style_spec(master_path).to_json())

    assert set(document) >= {
        "theme", "brand", "layouts", "grid", "chart_style",
        "table_style", "icon_palette", "meta",
    }
    assert set(document["theme"]) == {"colors", "fonts"}
    assert set(document["brand"]) == {"logo", "footer"}
    assert {"spec_version", "client", "project", "extracted_at"} <= set(document["meta"])
    assert {"major_latin", "major_ea", "major_cs"} <= set(document["theme"]["fonts"])
    assert {"source", "archetype", "placeholders"} <= set(document["layouts"][0])
