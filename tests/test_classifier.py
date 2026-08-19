"""Stage 3a: reading a rough slide for what it structurally is."""
from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from pptx_formatter import archetypes
from pptx_formatter.classifier import (
    assign_roles, classify_slide, extract_features, identify_title,
)

from conftest import textbox


def classify(prs, slide, index=0, total=1):
    return classify_slide(slide, index, total, prs.slide_width, prs.slide_height)


def test_cover_slide(widescreen):
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, "Q3 Growth Strategy", 1, 2.6, 9, 1.3, size=40)
    textbox(slide, "Prepared for the board", 1, 4.1, 9, 0.6, size=16)
    assert classify(widescreen, slide, index=0).archetype == archetypes.TITLE_SLIDE


def test_title_and_content(widescreen):
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, "Market Overview", 0.8, 0.5, 10, 0.9, size=32)
    textbox(slide, "Demand grew 14%\nCompetitors slowed\nMargins held", 0.8, 1.9, 10, 3, size=18)
    assert classify(widescreen, slide, index=2).archetype == archetypes.TITLE_AND_CONTENT


def test_two_column_from_horizontal_split(widescreen):
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, "Where We Win", 0.8, 0.5, 11, 0.9, size=32)
    textbox(slide, "Enterprise renews at 94% and keeps expanding steadily.",
            0.8, 2.1, 5.2, 2.6, size=18)
    textbox(slide, "Mid-market grows fastest but churn remains stubborn.",
            7.0, 2.1, 5.2, 2.6, size=18)
    assert classify(widescreen, slide, index=2).archetype == archetypes.TWO_CONTENT


def test_comparison_from_paired_headings(widescreen):
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, "Build or Partner", 0.8, 0.5, 11, 0.9, size=32)
    textbox(slide, "Build", 0.8, 1.9, 5.2, 0.5, size=20)
    textbox(slide, "Full control of the roadmap but fourteen months to parity.",
            0.8, 2.6, 5.2, 2.2, size=16)
    textbox(slide, "Partner", 7.0, 1.9, 5.2, 0.5, size=20)
    textbox(slide, "Live within a quarter but a twenty-two percent revenue share.",
            7.0, 2.6, 5.2, 2.2, size=16)
    assert classify(widescreen, slide, index=2).archetype == archetypes.COMPARISON


def test_chart_beats_everything(widescreen):
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, "Revenue by Quarter", 0.8, 0.5, 10, 0.9, size=32)
    data = CategoryChartData()
    data.categories = ["Q1", "Q2"]
    data.add_series("Revenue", (1.0, 2.0))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2), Inches(8), Inches(4), data
    )
    result = classify(widescreen, slide, index=2)
    assert result.archetype == archetypes.CHART
    assert result.confidence >= 0.9


def test_table_is_detected(widescreen):
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, "Key Accounts", 0.8, 0.5, 10, 0.9, size=32)
    slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(8), Inches(2))
    assert classify(widescreen, slide, index=2).archetype == archetypes.TABLE


def test_quote_from_punctuation_and_attribution(widescreen):
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, '"This reset how our team works."\n— Chief Executive, Northwind',
            1.6, 2.6, 9.5, 2.2, size=26)
    assert classify(widescreen, slide, index=3).archetype == archetypes.QUOTE


def test_statistic_slide(widescreen):
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, "94% net revenue retention", 1.6, 3.0, 10, 1.6, size=48)
    assert classify(widescreen, slide, index=4).archetype == archetypes.BIG_STATEMENT


def test_closing_slide_in_english_and_arabic(widescreen):
    for text in ("Thank you", "شكرا"):
        slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
        textbox(slide, text, 4.6, 3.2, 5, 1.4, size=44)
        assert classify(widescreen, slide, index=9, total=10).archetype == archetypes.CLOSING


def test_blank_slide(widescreen):
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    assert classify(widescreen, slide).archetype == archetypes.BLANK


def test_arabic_slide_is_flagged_rtl(widescreen):
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, "نظرة عامة على السوق", 0.8, 0.5, 10, 0.9, size=32)
    textbox(slide, "نمت المبيعات بنسبة 14% مقارنة بالعام الماضي", 0.8, 1.9, 10, 3, size=18)
    result = classify(widescreen, slide, index=2)
    assert result.features.is_rtl


def test_classification_carries_its_evidence(widescreen):
    """A designer overruling the tool should be able to see why it chose."""
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, "Key Accounts", 0.8, 0.5, 10, 0.9, size=32)
    slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(8), Inches(2))
    assert any("table" in line for line in classify(widescreen, slide).evidence)


def test_title_promoted_on_title_dominant_layouts(widescreen):
    """
    A rough "Thank you" slide has one text box and no title placeholder.
    Positionally it reads as body copy; knowing it's a closing says otherwise.
    """
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, "Thank you", 4.6, 3.2, 5, 1.4, size=44)
    features = extract_features(slide, 9, 10, widescreen.slide_width, widescreen.slide_height)

    assert identify_title(features) is None       # too low on the slide
    title, bodies = assign_roles(features, archetypes.CLOSING)
    assert title is not None
    assert not bodies


def test_quote_splits_from_its_attribution(widescreen):
    """The quotation and the credit belong in different placeholders."""
    slide = widescreen.slides.add_slide(widescreen.slide_layouts[6])
    textbox(slide, '"This reset how we work."\n— Chief Executive', 1.6, 2.6, 9.5, 2.2, size=26)
    features = extract_features(slide, 3, 6, widescreen.slide_width, widescreen.slide_height)

    title, bodies = assign_roles(features, archetypes.QUOTE)
    assert title is not None and len(bodies) == 1
    assert "reset how we work" in "".join(p.text for p in title.para_list())
    assert "Chief Executive" in "".join(p.text for p in bodies[0].para_list())
