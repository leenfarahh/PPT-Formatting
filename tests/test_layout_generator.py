"""Stage 2: building a full layout set on the designer's own master."""
from __future__ import annotations

from lxml import etree
from pptx import Presentation

from pptx_formatter import archetypes
from pptx_formatter.extraction import extract_style_spec
from pptx_formatter.layout_builder import add_layout, generated_layout_spec
from pptx_formatter.layout_generator import (
    generate_master_layouts, normalize_theme_references, strip_content_slides,
)
from pptx_formatter.style_spec import SOURCE_DESIGNER, StyleSpec

from conftest import assert_valid_pptx

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def test_designer_layouts_survive(master_path):
    """The point of building on their master: their layouts are kept as-is."""
    spec = extract_style_spec(master_path)
    original_names = {l.name for l in spec.layouts}

    prs, report = generate_master_layouts(spec, master_path)
    produced = set(report["available_layouts"])

    assert original_names <= produced
    assert len(report["designer_layouts"]) == len(spec.archetypes_present())


def test_content_slides_are_stripped(master_path):
    spec = extract_style_spec(master_path)
    prs, report = generate_master_layouts(spec, master_path)
    assert report["slides_stripped"] == 2
    assert len(prs.slides._sldIdLst) == 0


def test_missing_archetypes_are_filled(master_path):
    """Gaps the designer left get real layouts, not approximations."""
    spec = extract_style_spec(master_path)
    prs, report = generate_master_layouts(spec, master_path)

    covered = set(report["layout_archetypes"].values())
    for archetype in archetypes.ALL_ARCHETYPES:
        assert archetype in covered, f"no layout covers {archetype}"

    filled = {item["archetype"] for item in report["generated_layouts"]}
    assert archetypes.QUOTE in filled
    assert archetypes.CHART in filled


def test_generated_layouts_sit_on_the_designers_grid(master_path):
    """A synthesized layout should respect their margins, not a default."""
    spec = extract_style_spec(master_path)
    layout = generated_layout_spec(archetypes.TITLE_AND_CONTENT, spec)

    title = next(p for p in layout.placeholders if p.is_title)
    assert abs(title.left_frac - spec.grid.margin_left_frac) < 1e-6
    assert abs(title.top_frac - spec.grid.margin_top_frac) < 1e-6


def test_added_layout_is_usable_and_survives_a_round_trip(master_path, tmp_path):
    """
    python-pptx can't add layouts, so this exercises the raw part wiring:
    the new layout has to be real enough to add a slide from and to reopen.
    """
    spec = extract_style_spec(master_path)
    prs, _ = generate_master_layouts(spec, master_path)

    quote = next(l for l in prs.slide_masters[0].slide_layouts if l.name == "Quote")
    slide = prs.slides.add_slide(quote)
    slide.shapes.title.text = "A synthesized layout"

    out = tmp_path / "layouts.pptx"
    prs.save(str(out))
    assert_valid_pptx(out)

    reopened = Presentation(str(out))
    assert "Quote" in [l.name for l in reopened.slide_masters[0].slide_layouts]
    assert reopened.slides[0].slide_layout.name == "Quote"
    assert reopened.slides[0].shapes.title.text == "A synthesized layout"


def test_theme_is_restyled(master_path):
    """Every layout follows the theme, so rewriting it restyles all of them."""
    spec = extract_style_spec(master_path)
    spec.theme.colors["accent1"] = "FF0000"
    spec.theme.fonts.minor_latin = "Verdana"

    prs, _ = generate_master_layouts(spec, master_path)
    theme_root = etree.fromstring(
        prs.slide_masters[0].part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        ).blob
    )
    accent1 = theme_root.find(f".//{{{A}}}clrScheme/{{{A}}}accent1/{{{A}}}srgbClr")
    minor = theme_root.find(f".//{{{A}}}fontScheme/{{{A}}}minorFont/{{{A}}}latin")
    assert accent1.get("val") == "FF0000"
    assert minor.get("typeface") == "Verdana"


def test_complex_script_font_is_always_written(master_path):
    """Formatting points Arabic at +mn-cs, so that entry has to resolve."""
    spec = extract_style_spec(master_path)
    spec.theme.fonts.major_cs = ""
    spec.theme.fonts.minor_cs = ""
    spec.theme.fonts.minor_latin = "Verdana"

    prs, _ = generate_master_layouts(spec, master_path)
    theme_root = etree.fromstring(
        prs.slide_masters[0].part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        ).blob
    )
    cs = theme_root.find(f".//{{{A}}}fontScheme/{{{A}}}minorFont/{{{A}}}cs")
    assert cs.get("typeface") == "Verdana", "cs must fall back rather than stay empty"


def test_hardcoded_values_become_theme_references():
    """
    The inheritance repair: a layout pinning a theme color and font looks
    right today but stops following the master the moment the theme changes.
    """
    spec = StyleSpec()
    spec.theme.colors.update({"accent1": "0F4C81", "dk1": "1A1A2E"})
    spec.theme.fonts.minor_latin = "Segoe UI"

    element = etree.fromstring(
        f'<root xmlns:a="{A}">'
        f'  <a:latin typeface="Segoe UI"/>'
        f'  <a:solidFill><a:srgbClr val="0F4C81"><a:alpha val="50000"/></a:srgbClr></a:solidFill>'
        f'  <a:srgbClr val="ABCDEF"/>'
        f"</root>".encode()
    )
    changes = normalize_theme_references(element, spec)

    assert element.find(f"{{{A}}}latin").get("typeface") == "+mn-lt"
    scheme_clr = element.find(f".//{{{A}}}schemeClr")
    assert scheme_clr.get("val") == "accent1"
    # Transforms on the original color must survive the swap.
    assert scheme_clr.find(f"{{{A}}}alpha") is not None
    # A color that isn't in the theme is a design decision; leave it alone.
    assert element.find(f"{{{A}}}srgbClr").get("val") == "ABCDEF"
    assert len(changes) == 2


def test_dark_and_light_roles_map_through_the_color_map():
    spec = StyleSpec()
    spec.theme.colors.update({"dk1": "000000", "lt1": "FFFFFF"})
    element = etree.fromstring(
        f'<root xmlns:a="{A}"><a:srgbClr val="000000"/><a:srgbClr val="FFFFFF"/></root>'.encode()
    )
    normalize_theme_references(element, spec)
    values = [el.get("val") for el in element.findall(f"{{{A}}}schemeClr")]
    assert values == ["tx1", "bg1"]


def test_logo_is_placed_on_the_master(master_path, tmp_path, widescreen):
    """One copy on the master, inherited everywhere, rather than per layout."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from examples.make_sample_master import write_solid_png

    logo = tmp_path / "logo.png"
    write_solid_png(logo, 120, 40, (15, 76, 129))

    spec = extract_style_spec(master_path)
    spec.brand.logo.asset_path = str(logo)
    spec.brand.logo.left_frac, spec.brand.logo.top_frac = 0.85, 0.05
    spec.brand.logo.width_frac, spec.brand.logo.height_frac = 0.11, 0.07

    prs, _ = generate_master_layouts(spec, master_path)
    pictures = [
        s for s in prs.slide_masters[0].shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert len(pictures) == 1

    # Re-running must not stack a second copy.
    prs2, _ = generate_master_layouts(spec, master_path)
    assert len([
        s for s in prs2.slide_masters[0].shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]) == 1


def test_bank_fills_gaps_before_generating(master_path, bank, tmp_path):
    """A banked layout should be preferred over a synthesized one."""
    from pptx_formatter.pipeline import ingest_master

    # First client contributes a quote layout to the bank.
    donor = extract_style_spec(master_path, client="Donor")
    quote = generated_layout_spec(archetypes.QUOTE, donor)
    quote.name = "Donor Quote"
    donor.layouts.append(quote)
    bank.save(donor, "donor", master_pptx=master_path)

    # Second client's master has no quote layout of its own.
    spec = extract_style_spec(master_path, client="Second")
    assert archetypes.QUOTE not in spec.archetypes_present()

    prs, report = generate_master_layouts(spec, master_path, bank=bank)
    sourced = {item["archetype"]: item["source"] for item in report["bank_layouts"]}
    assert sourced.get(archetypes.QUOTE) == "bank:donor"
