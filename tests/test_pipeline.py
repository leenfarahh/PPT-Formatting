"""End-to-end: master plus rough deck becomes a formatted, editable deck."""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from pptx_formatter import archetypes
from pptx_formatter.pipeline import (
    format_with_banked_spec, ingest_master, run_pipeline,
)

from conftest import assert_valid_pptx, textbox


def test_full_pipeline(master_path, content_path, tmp_path, bank):
    out = tmp_path / "formatted.pptx"
    report = run_pipeline(
        master_path, content_path, out, bank=bank,
        client="Acme Holdings", project="Board Deck",
    )

    assert out.exists()
    assert_valid_pptx(out)
    assert report["slides_processed"] == 6
    assert report["master"]["content_slides_ignored"] == 2


def test_output_opens_and_is_placeholder_backed(master_path, content_path, tmp_path, bank):
    """
    Content mapped into placeholders is what makes the deck editable and
    reflowable rather than a pile of pinned text boxes.
    """
    out = tmp_path / "formatted.pptx"
    run_pipeline(master_path, content_path, out, bank=bank, client="Acme")

    reopened = Presentation(str(out))
    first = reopened.slides[0]
    assert first.shapes.title is not None
    assert "Q3 Growth Strategy" in first.shapes.title.text


def test_every_slide_lands_on_a_matching_layout(master_path, content_path, tmp_path, bank):
    report = run_pipeline(master_path, content_path, out := tmp_path / "o.pptx", bank=bank)
    routed = {s["slide"]: s["archetype"] for s in report["slides"]}

    assert routed[1] == archetypes.TITLE_SLIDE
    assert routed[3] == archetypes.TWO_CONTENT
    assert routed[5] == archetypes.QUOTE
    assert routed[6] == archetypes.CLOSING
    # Nothing may fall back to an arbitrary layout.
    assert all(s["layout"] for s in report["slides"])


def test_layout_can_be_overridden(master_path, content_path, tmp_path, bank):
    """A designer disagreeing with the classifier shouldn't have to edit code."""
    report = run_pipeline(
        master_path, content_path, tmp_path / "o.pptx", bank=bank,
        layout_overrides={2: archetypes.SECTION_HEADER},
    )
    slide = report["slides"][2]
    assert slide["archetype"] == archetypes.SECTION_HEADER
    assert any("overridden" in line for line in slide["evidence"])


def test_arabic_slide_gets_rtl_treatment(master_path, content_path, tmp_path, bank):
    report = run_pipeline(master_path, content_path, tmp_path / "o.pptx", bank=bank)
    arabic = report["slides"][3]
    assert arabic["formatting"]["rtl_paragraphs"] >= 1


def test_charts_survive_with_their_data(master_path, tmp_path, bank):
    """
    A chart points at its own part plus an embedded workbook. Copying only
    the shape XML leaves those dangling and the file won't open.
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox(slide, "Revenue by Quarter", 0.8, 0.5, 10, 0.9, size=32)
    data = CategoryChartData()
    data.categories = ["Q1", "Q2", "Q3"]
    data.add_series("Enterprise", (18.2, 22.5, 26.1))
    data.add_series("Mid-market", (9.4, 10.1, 12.8))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2), Inches(8), Inches(4), data
    )
    content = tmp_path / "charts.pptx"
    prs.save(str(content))

    out = tmp_path / "formatted.pptx"
    report = run_pipeline(master_path, content, out, bank=bank)
    assert_valid_pptx(out)

    reopened = Presentation(str(out))
    charts = [
        s for s in reopened.slides[0].shapes
        if getattr(s, "has_chart", False) and s.has_chart
    ]
    assert len(charts) == 1
    chart = charts[0].chart
    assert len(list(chart.series)) == 2
    assert list(chart.plots[0].categories) == ["Q1", "Q2", "Q3"]
    assert report["slides"][0]["formatting"]["charts_formatted"] == 1


def test_tables_survive(master_path, tmp_path, bank):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox(slide, "Key Accounts", 0.8, 0.5, 10, 0.9, size=32)
    table = slide.shapes.add_table(
        3, 3, Inches(1), Inches(2), Inches(8), Inches(2)
    ).table
    for r, row in enumerate([["A", "B", "C"], ["1", "2", "3"], ["4", "5", "6"]]):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
    content = tmp_path / "tables.pptx"
    prs.save(str(content))

    out = tmp_path / "formatted.pptx"
    run_pipeline(master_path, content, out, bank=bank)
    assert_valid_pptx(out)

    reopened = Presentation(str(out))
    tables = [
        s for s in reopened.slides[0].shapes
        if getattr(s, "has_table", False) and s.has_table
    ]
    assert len(tables) == 1
    assert tables[0].table.cell(0, 0).text == "A"


def test_no_content_is_silently_dropped(master_path, content_path, tmp_path, bank):
    """Anything that can't be mapped must be reported, not quietly lost."""
    report = run_pipeline(master_path, content_path, tmp_path / "o.pptx", bank=bank)
    for slide in report["slides"]:
        assert slide["mapped"] or slide["carried_over"] or slide["warnings"], (
            f"slide {slide['slide']} produced no account of its content"
        )


def test_repeat_client_skips_stage_one(master_path, content_path, tmp_path, bank):
    """The payoff of banking: a second deck needs no master submission."""
    ingest_master(master_path, bank=bank, client="Acme Holdings")

    out = tmp_path / "repeat.pptx"
    report = format_with_banked_spec("Acme Holdings", content_path, out, bank)

    assert report["stage_1_skipped"] is True
    assert report["slides_processed"] == 6
    assert_valid_pptx(out)


def test_unknown_client_is_a_clear_error(content_path, tmp_path, bank):
    with pytest.raises(LookupError, match="No banked Style Spec"):
        format_with_banked_spec("Nobody", content_path, tmp_path / "o.pptx", bank)


def test_pipeline_runs_without_a_bank(master_path, content_path, tmp_path):
    """The bank is an optimization, not a dependency."""
    out = tmp_path / "formatted.pptx"
    report = run_pipeline(master_path, content_path, out, bank=None)
    assert report["bank_entry"] is None
    assert report["slides_processed"] == 6
    assert_valid_pptx(out)


def test_empty_master_still_produces_a_full_layout_set(tmp_path, content_path, bank):
    """A designer may submit a bare master; every archetype still gets built."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    master = tmp_path / "empty.pptx"
    prs.save(str(master))

    out = tmp_path / "formatted.pptx"
    report = run_pipeline(master, content_path, out, bank=bank)

    covered = {s["archetype"] for s in report["slides"]}
    assert covered
    assert_valid_pptx(out)


def test_report_explains_where_each_layout_came_from(master_path, content_path, tmp_path, bank):
    report = run_pipeline(master_path, content_path, tmp_path / "o.pptx", bank=bank)
    layouts = report["layouts"]

    assert layouts["designer"], "designer-authored layouts should be reported"
    assert layouts["generated"], "gaps should be reported as generated"
    for item in layouts["generated"]:
        assert item["source"] == "generated"
    assert "inheritance" in layouts
